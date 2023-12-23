# Copyright (c) 2020 Neteera Technologies Ltd. - Confidential

import os
import sys

sys.path.append(os.path.dirname(os.path.dirname(os.path.dirname(__file__))))    # noqa

from Configurations import back_chair_config
from Offline import create_config
from Spot import standard_hr_filter, Radar, normal_round
from SignalBuffer import SignalBuffer
from ChannelSelectors import fmcw_channel_select, compute_dr
from pylibneteera.filters import Filter

from pylibneteera.float_indexed_array import TimeArray
from pylibneteera.float_indexed_array import FrequencyArray
from pylibneteera.sp_utils import spectrogram, normalize

from Tests.Utils.IOUtils import load
from Tests.Utils.ResearchUtils import plot
from Tests.Utils.DBUtils import find_back_setup_same_session, match_lists_by_ts
from Tests.Utils.LoadingAPI import load_nes, load_reference, load_data_from_json
from Tests.vsms_db_api import *
from Tests.Utils.PathUtils import change_dir_to_pythondev

import argparse
import datetime
from io import BytesIO
import matplotlib.pyplot as plt
from matplotlib.ticker import AutoMinorLocator
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.util import Inches
from tkinter import filedialog
import tkinter as tk
from scipy.signal import welch
import seaborn as sbs

from bokeh.io import show
from bokeh.plotting import figure, output_file, save
from bokeh.models import ColumnDataSource, Circle
from bokeh.layouts import gridplot

C = 3e8


def get_args() -> argparse.Namespace:
    """ Argument parser

    :return: parsed arguments of the types listed within the function
    :rtype: argparse.Namespace
    """
    parser = argparse.ArgumentParser(description='Process some integers.')
    parser.add_argument('-setups', '-session_ids', '-setup_ids', '-setup',
                        metavar='ids', nargs='+', type=int, help='Setup IDs in DB', required=False)
    parser.add_argument('-tlog_paths', metavar='LoadPath', nargs='+', type=str, help='Path to tlogs', required=False)
    parser.add_argument('-json_paths', metavar='LoadPath', nargs='+', type=str, help='path to jsons', required=False)
    parser.add_argument('-save_path', metavar='SavePath', type=str, help='Path to save images', required=False)
    parser.add_argument('--dont_track', action='store_true', help='Track the offset', required=False)
    parser.add_argument('--dont_overwrite', action='store_true', help='Dont overwrite existing images', required=False)
    parser.add_argument('--dont_show', action='store_true', help='Dont display the created images', required=False)
    parser.add_argument('--plot_status', action='store_true', help='Display ref recorder on displacement plot',
                        required=False)
    parser.add_argument('-ppt', type=str, required=False, help='Ppt filename to generate presentation from plots')
    parser.add_argument('--nlf', required=False, action='store_true', help='Non-linear filter')
    parser.add_argument('-plots', metavar='plots', nargs='+', type=str, help='Plots to create', required=False,
                        default=['iq', 'freqs', 'freqs_nlf', 'psd', 'iq_vs_time', 'displacement', 'amplitude'],
                        choices=['iq', 'freqs', 'freqs_nlf', 'psd', 'iq_vs_time', 'displacement', 'amplitude',
                                 'bokeh', 'displacement_lpf'])
    parser.add_argument('-select_range', metavar='dividers', nargs=2, type=int, help='select bin range', required=False)
    parser.add_argument('--variance', action='store_true',
                        help='if variance will be showed, other plots will not be shown', required=False)
    parser.add_argument('--back_twin', action='store_true',
                        help='plot a back setup recorded in this setup as well', required=False)
    parser.add_argument('-time_range', metavar='dividers', nargs=2, type=int,  help='select time_vec range for iq plot',
                        required=False)
    parser.add_argument('-intra_breath_path', type=str, help='path to intra_breath results', required=False)
    parser.add_argument('-distance', type=int, help='distance from radar to subject in mm', required=False)
    parser.add_argument('--update_offset_dynamically', action='store_true',
                        help='update offset every 20 seconds (default is calc only for the first seconds)')
    parser.add_argument('-mysql_db', '-db', '-db_name', type=str, default='neteera_db',
                        help='choose mysql server: neteera_db - for setups recorded with tlog neteera_cloud_mirror for'
                             'setups from the cloud',
                        choices=['neteera_db', 'neteera_cloud_mirror'])
    return parser.parse_args()


def bins_select(params: dict, dr) -> dict:
    """ Selects the bins to plot"""
    args = get_args()
    bin_range = args.select_range
    if bin_range[1] <= 40:
        bin_range = dr * np.array(bin_range)
    bins = {b: dr * b for b in params['bins'] if bin_range[0] <= dr * b <= bin_range[1]}
    return bins


def get_distance():
    master = tk.Tk()
    master.wm_attributes('-topmost', 1)
    tk.Label(master, text="Distance [mm]").grid(row=0)
    e1 = tk.Entry(master)
    e1.grid(row=0, column=1)
    tk.Button(master, text='Go', command=master.quit).grid(row=3, column=0, sticky=tk.W, pady=4)
    tk.mainloop()
    s = e1.get()
    master.destroy()
    if all([ss.isdigit() for ss in s.split('.')]) and len(s) > 0:
        return float(s)


def load_data_from_tlog(file_path, db, dist=None):
    print('loading data from {}'.format(file_path))
    data = load_nes(path=file_path, db=db)
    data['gap'] = 1 / data['framerate']
    if dist is None:
        data['dist'] = get_distance()
    else:
        data['dist'] = dist
    return data


def load_data_from_db(idx, db):
    print(f'loading data from db setup {idx}')
    data = load_nes(idx, db)
    if len(data) == 4:
        return data
    db_dist = db.setup_distance(idx)
    if db_dist:
        data['dist'] = db_dist
    data['gap'] = 1 / data['framerate']
    return data


def roll_data(readings, channel_select_params, gap, win_size, dont_track, update_offset_dynamic=False, bin_index=None):
    if bin_index is None:
        bin_index = fmcw_channel_select(readings.iloc[:int(10 / gap), :], channel_select_params, back_chair_config)
    if dont_track:
        return readings.iloc[:, bin_index]
    sb = SignalBuffer(back_chair_config)
    sb.bin_index = bin_index
    output = []
    readings_batch_size = int(win_size / gap)
    for i in range(0, len(readings), readings_batch_size):
        sb.push(TimeArray(readings[i: i + readings_batch_size], gap=gap), channel_select_params)
        if update_offset_dynamic or i == 0:
            sb.update_offset(win_size)
        output.append(sb.raw_window(win_size, True))
    return np.concatenate(output)


def gen_label(dont_track, setup_id, path, db):
    str_tracked = 'offset not tracked' if dont_track else 'offset tracked'
    if setup_id is not None:
        target = db.setup_target(setup_id)
        distance = db.setup_distance(setup_id)
        return f'setup {setup_id} Target {target} offset {str_tracked} tracked distance {distance}mm' \
               f'\n {db.setup_sn(setup_id)}'
    else:
        return f'{os.path.basename(path)} {str_tracked}'


def gen_save_path(save_path, setu_id=None, path=None):
    if save_path is not None:
        if setu_id is None:
            base_name = os.path.basename(path)[:os.path.basename(path).find('_')] + '_'
        else:
            base_name = f'{setu_id}_'
        return os.path.join(save_path, base_name)


def gen_disp_time(phases, gap):
    """See /Neteera/Work/Departments/Algo/Algo_and_Radar_course/CW radar principles_part1.mp4 10:00-17:00
    :return time_vector and displacement in micro-meters"""
    # convert to microns
    wave_length = 3e8 / 124e9
    k = 2 * np.pi / wave_length
    displacement = phases / (2 * k)
    time_vec = np.arange(len(displacement)) * gap
    return time_vec, displacement * 1e6


def images_exists(cur_save_path, dont_track):
    if dont_track:
        if os.path.isfile(cur_save_path + 'iq_offset_not_tracked.svg') and os.path.isfile(
                cur_save_path + 'displacement_offset_not_tracked.svg'):
            return True
    else:
        if os.path.isfile(cur_save_path + 'iq_offset_tracked.svg') and os.path.isfile(
                cur_save_path + 'displacement_offset_tracked.svg'):
            return True
    return False


def win_ft(signal, start, win_sec, fs, bp_band, padding):
    window_raw = signal[start * fs: (start + win_sec) * fs]
    window = Filter().butter_bandpass_filter(TimeArray(window_raw, 1/fs), *bp_band)
    nfft = int(2 ** (round(np.log2(win_sec * fs)) + padding))
    return spectrogram(window, nfft)


def signal_ft(signal, win_sec, fs, bp_band, padding):
    nfft = int(2 ** (round(np.log2(win_sec * fs)) + padding))
    return FrequencyArray([win_ft(signal, ii, win_sec, fs, bp_band, padding) for ii in range(0, len(signal) // fs, 1)],
                          fs / nfft)


def _fig_visualization(name, grid):
    if len(plt.gca().get_title()) > 0:
        plt.title(name)
    plt.grid(grid)
    plt.legend()


def _save_png_and_svg(save_path):
    if save_path is not None:
        png_save_path = save_path
        svg_save_path = png_save_path.replace('.png', '.svg')
        plt.savefig(png_save_path)
        print(f'\nplot saved to {png_save_path}\n')
        plt.savefig(svg_save_path)
        print(f'\nplot saved to {svg_save_path}\n')


def _add_fig_to_ppt():
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    image_stream = BytesIO()
    plt.savefig(image_stream)
    slide.shapes.add_picture(image_stream, Inches(2.5), Inches(1.5), height=Inches(3.1))


def _plot_closer(full_label, save_path, dont_show, ppt, grid=True):
    _fig_visualization(full_label, grid)
    _save_png_and_svg(save_path)
    if ppt:
        _add_fig_to_ppt()
    if dont_show:
        plt.close()


def plot_bokeh(data, label, save_path, bin_fmcw, displacement, amplitude, time_vec):
    data_dict = data[bin_fmcw]
    i = data_dict['i']
    q = data_dict['q']

    points_number = 1500
    i_sub_sampled = i[0: -1: 1 + int(len(i) / points_number)]
    q_sub_sampled = q[0: -1: 1 + int(len(i) / points_number)]
    sub_time = time_vec[0: -1: 1 + int(len(i) / points_number)]

    if save_path is not None:
        output_file(f'{save_path}{label}_bin{bin_fmcw}_bokeh.html')

    # create a column data source for the plots to share
    source = ColumnDataSource(data=dict(i_sub=i_sub_sampled, q_sub=q_sub_sampled, t_sub=sub_time,
                                        disp=displacement[0: -1: 1 + int(len(i) / points_number)],
                                        amp=amplitude[0: -1: 1 + int(len(i) / points_number)]))

    TOOLS = "box_select,help, reset, box_zoom, hover, pan, undo"

    # create a new plot and add a renderer
    iq = figure(tools=TOOLS, width=550, height=550, title=label)
    render = iq.circle('i_sub', 'q_sub', source=source)
    selected_circle = Circle(fill_color='firebrick', size=10)
    render.selection_glyph = selected_circle
    iq.xaxis.axis_label = 'In Phase'
    iq.yaxis.axis_label = 'Quadrature'

    # create another new plot and add a renderer
    iq_vs_time = figure(tools=TOOLS, width=550, height=400, title='iq_vs_time', toolbar_location='right')
    iq_vs_time.circle('t_sub', 'i_sub', source=source, fill_alpha=0, line_color=None)
    iq_vs_time.line('t_sub', 'i_sub', source=source, legend_label='I')
    iq_vs_time.line('t_sub', 'q_sub', source=source, line_color='#ff7f0e', legend_label='Q')
    iq_vs_time.xaxis.axis_label = 'time_vec [sec]'

    disp_fig = figure(tools=TOOLS, width=550, height=400, title='displacement', x_range=iq_vs_time.x_range)
    disp_fig.line('t_sub', 'disp', source=source)
    disp_fig.xaxis.axis_label = 'time_vec [sec]'
    disp_fig.yaxis.axis_label = 'displacement [microns]'

    amp_fig = figure(tools=TOOLS, width=550, height=400, title='|I + jQ|', x_range=iq_vs_time.x_range)
    amp_fig.line('t_sub', 'amp', source=source)
    amp_fig.xaxis.axis_label = 'time_vec [sec]'

    p = gridplot([[iq, iq_vs_time], [amp_fig, disp_fig]], toolbar_location='right')

    if save_path is not None:
        save(p)
    else:
        show(p)


def plot_iq(data, label, save_path, ppt, dont_show, bin_fmcw, time_range, first_bin, is_multi):
    data_dict = data[bin_fmcw]
    i = data_dict['i']
    q = data_dict['q']
    fs = 1 / data_dict['gap']
    if time_range is None:
        start = 0
        end = len(data_dict['i']) / fs
    else:
        start = time_range[0]
        end = time_range[1]

    i_sub_sampled = i[0: -1: 1 + int(len(i) / 15000)]
    q_sub_sampled = q[0: -1: 1 + int(len(i) / 15000)]
    time_vec = np.linspace(start, end, len(i_sub_sampled))

    full_label = 'I_vs_Q ' + label
    plt.figure(full_label)
    if is_multi:
        plt.scatter(i_sub_sampled, q_sub_sampled, label=bin_fmcw, s=1)
    else:
        plt.scatter(i_sub_sampled, q_sub_sampled, label=bin_fmcw, c=time_vec, cmap='tab20', s=1)
        if first_bin:
            plt.colorbar()
    plt.xlabel(r'$In\ Phase$')
    plt.ylabel(r'$Quadrature$')
    limit = np.max(np.abs(np.concatenate((i_sub_sampled, q_sub_sampled))))
    limit = max(abs(plt.xlim()[0]), np.ceil(limit * 1.1))
    plt.xlim([-limit, limit])
    plt.ylim([-limit, limit])

    _plot_closer(full_label, save_path, dont_show, ppt)


def plot_iq_vs_time(time_vec, i, q, label, save_path, ppt, dont_show, bin_fmcw=None):
    full_label = 'iq_vs_time ' + label
    plt.figure(full_label)
    plt.plot(time_vec, i[(i.shape[0] - time_vec.shape[0]):], label=f"I {bin_fmcw}", linewidth=0.7)
    plt.plot(time_vec, q[(q.shape[0] - time_vec.shape[0]):], label=f"Q {bin_fmcw}", linewidth=0.7)
    plt.xlabel(r'$Time\ [seconds]$')
    plt.ylabel(r'Counts')
    _plot_closer(full_label, save_path, dont_show, ppt)


def plot_amplitude(timee, amplitude, label, save_path, ppt, dont_show, bin_fmcw=None):
    full_label = 'Amplitude ' + label
    plt.figure(full_label)
    plt.plot(timee, amplitude, label=bin_fmcw, linewidth=0.7)
    plt.xlabel(r'$Time\ [seconds]$')
    plt.ylabel(r'Amplitude')

    mean_div_perc_vec = []
    for i in range(0, len(amplitude) - 500*20, 500):
        window = amplitude[i: i + 20 * 500]
        mean_div_perc_vec.append(np.mean(window) / np.percentile(window, 5))

    _plot_closer(full_label, save_path, dont_show, ppt)


def plot_status(idx, timee):
    time_nes = np.arange(int(timee[0]), int(timee[-1] + 1.5))
    ref_motion = load_reference(idx, 'motion', db)  # match time_vec vecs once, on motion with no special reasons
    if ref_motion is None:
        return
    time_ref = np.arange(len(ref_motion))
    new_t_nes, new_t_ref, _ = match_lists_by_ts(time_nes, time_ref[int(timee[0]):], idx, 'rest', db)

    for status in ['motion', 'speaking', 'zrr', 'occupancy']:
        ref_status = np.array(load_reference(idx, status, db, True))
        try:
            ref = ref_status[new_t_ref]
        except IndexError:
            continue
        if status == 'occupancy':
            ref = 1 - ref
            status = 'empty chair'
        if len(ref[ref == 1]):
            plt.plot(new_t_nes[ref == 1], ref[ref == 1], '.', label=status)
    plt.legend()


def plot_intra_breath_indicator(timee, displacement, idx, folder):
    timee_downsampled = timee[::50]
    displacement_downsampled = displacement[::50]
    indicator = load(os.path.join(folder, f'{idx}_intra_breath.npy'), allow_pickle=True)
    indicator = np.array([1 if x == True else np.nan for x in indicator])[:-10] # noqa
    plt.plot(timee_downsampled, indicator * displacement_downsampled, 'o', label='indicator', c='g')


def plot_displacement(timee, displacement, label, save_path, ppt, dont_show, st_plot, idx=None, bin_fmcw=None,
                      intra_breath_path=None):
    full_label = 'Displacement ' + label
    plt.figure(full_label)
    plt.plot(timee, displacement, label=bin_fmcw, linewidth=0.7)
    plt.xlabel(r'$Time\ [seconds]$')
    plt.ylabel(r'$Displacement\ [\mu m]$')

    if idx:
        plot_status(idx, timee)
    if intra_breath_path is not None:
        plot_intra_breath_indicator(timee, displacement, idx, intra_breath_path)

    _plot_closer(full_label, save_path, dont_show, ppt)


def plot_freqs(freqs, band, reference_data, label, save_path, ppt, dont_show):
    full_label = 'Freqs ' + label
    plt.figure(full_label)
    admissible_freqs = freqs.T.freqs[np.logical_and(band[0] < freqs.T.freqs, band[1] > freqs.T.freqs)]
    admissible_fses = normalize(freqs.T[np.logical_and(band[0] < freqs.T.freqs, band[1] > freqs.T.freqs)], axis=0)
    ax = sbs.heatmap(admissible_fses, yticklabels=np.round(admissible_freqs * 60))
    yticks = ax.yaxis.get_major_ticks()
    for i in range(len(yticks)):
        if i % 20 != 1:
            yticks[i].set_visible(False)
    if reference_data is not None and not isinstance(reference_data, int):
        reference_data = pd.Series(reference_data).fillna(np.nan).to_numpy()
        plt.plot((reference_data[5:] - 60 * band[0]) / (60 * (band[1] - band[0]) / len(yticks)))
        plt.plot((np.asarray(reference_data[5:]) * 2 - 60 * band[0]) / (60 * (band[1] - band[0]) / len(yticks)))
    ax.invert_yaxis()
    plt.xlabel('time_vec [s]')
    plt.ylabel('BPM')
    _plot_closer(full_label, save_path, dont_show, ppt, grid=False)


def plot_psd(freqs_psd, power, reference_data, label, save_path, ppt, dont_show, bin_fmcw=None):
    full_label = 'psd' + label
    plt.figure(full_label)
    plt.xlim(left=0, right=500)
    plt.plot(freqs_psd * 60, power, label=bin_fmcw, linewidth=2)
    if reference_data is not None:
        plt.axvline(float(np.mean(reference_data)), color='orange')
        plt.axvline(float(np.mean(reference_data)) * 2, color='orange', linestyle='--')
    plt.xlabel('BPM')
    plt.ylabel('PSD [dB/Hz]')
    _plot_closer(full_label, save_path, dont_show, ppt)


def plot_variance(dr, channel_select_params, label, readings):
    variance_label = f'Variance of STD and Mean for bins in setup \n {label}'

    readings_t = readings.transpose()
    bin_mag = {}
    for i in range(len(readings_t)):
        bin_mag[i] = []
    for r in range(len(readings_t)):
        for c in range(len(readings_t[0])):
            bin_mag[r].append(np.absolute(readings_t[r][c]))
    std = []
    mean = []
    for k in bin_mag:
        std.append(np.std(bin_mag[k]))
        mean.append(np.mean(bin_mag[k]))
    fig, ax1 = plt.subplots()
    color = 'tab:red'
    ax1.plot(channel_select_params['bins'], std, color=color)
    plt.title(variance_label)
    plt.grid()
    ax1.set_xlabel('Bins')
    ax1.set_ylabel('STD of magnitude', color=color)
    ax1.tick_params(axis='y', labelcolor=color)

    color = 'tab:blue'
    ax2 = ax1.twinx()
    ax2.plot(channel_select_params['bins'], mean, color=color)
    ax2.set_ylabel('Mean of magnitude', color=color)  # we already handled the x-label with ax1
    ax2.tick_params(axis='y', labelcolor=color)

    def forward(x):
        return x * dr

    secax = ax1.secondary_xaxis('top', functions=(forward, forward))
    secax.xaxis.set_minor_locator(AutoMinorLocator())
    secax.set_xlabel('Distance [mm]')
    fig.tight_layout()
    plt.show()


def plot_raw(raw, label, save_path, dont_show, ppt):
    if raw is not None:
        full_label = 'Raw data ' + label
        plt.figure(full_label)
        interval = 10 ** (len(str(int(len(raw)))) - 2)
        for frame in range(0, interval * 10, interval):
            frame_sig = raw[frame]
            plot(frame_sig, label=f'frame_{frame}', fig_name=full_label)
        _plot_closer(full_label, save_path, dont_show, ppt)


def plot_bins(timee, bins, label, save_path, dont_show, ppt):
    # only works for data from cloud (json files)
    if bins is not None:
        full_label = 'bins ' + label
        plt.figure(full_label)
        time_start = int(timee[0])
        plt.plot(range(time_start, time_start + len(bins)), bins)
        plt.xlabel('Time [sec]')
        plt.ylabel('bin index')
        plt.ylim(0, 18)
        _plot_closer(full_label, save_path, dont_show, ppt)


def load_data(setup_id, tlog_path, json_path, distance, db):
    if setup_id is not None:
        return load_data_from_db(setup_id, db)
    elif tlog_path is not None:
        return load_data_from_tlog(tlog_path, db, distance)
    else:
        return load_data_from_json([json_path], db)


def cut_signal_by_time_range(readings, bin_nums, gap, base_label):
    readings = readings[int(args.time_range[0] / gap): int(args.time_range[1] / gap)]
    if len(readings) == 0:
        raise ValueError(f'no data in time range.\n args.time_range = {args.time_range} \n{base_label}')
    if bin_nums is not None:
        bin_nums = bin_nums[args.time_range[0]: args.time_range[1] + 1]
    return readings, bin_nums


def process_single_bin(bin_fmcw, dist, data, readings, gap, dont_track, update_offset_dynamically, setup_id, nlf, base_label,
                       bin_nums, bp_band, padding, psd_padding, overlap, win_size):
    data['dist'] = dist
    single_bin_data = roll_data(
        readings, data, gap, win_size, dont_track, update_offset_dynamically, bin_fmcw)
    i = np.real(single_bin_data)
    q = np.imag(single_bin_data)
    fs = int(1 / gap)
    amplitude = np.abs(single_bin_data)
    phases = np.unwrap(np.angle(single_bin_data), axis=0)
    time_vec, displacement = gen_disp_time(phases, gap)
    if args.time_range is not None:
        time_vec = time_vec + args.time_range[0]
    freqs = signal_ft(phases, win_size, fs, bp_band, padding)
    if db is not None:
        reference_data = load_reference(setup_id, 'hr', db)
        if args.time_range is not None and reference_data is not None:
            reference_data = reference_data[args.time_range[0]: args.time_range[1]]
    else:
        reference_data = None

    if nlf:
        r = Radar(
            {'i': np.real(i).reshape(-1),
             'q': np.imag(q).reshape(-1)
             },
            fs=fs,
            truncate=False
        )
        r.demodulate()
        phases_nlf = standard_hr_filter(r.data['complex_iq'],
                                        fs,
                                        high_f_band=(10, 20),
                                        low_f_band=[35. / 60., 600. / 60.])[0]
        freqs_nlf = signal_ft(phases_nlf, 10, fs, bp_band, padding)
        label = gen_label(dont_track, setup_id, base_label, db)
        plot_displacement(time_vec, phases_nlf, label + ' heart signal preprocessing', cur_save_path, ppt, False, False)
    else:
        freqs_nlf = None
    nfft = int(2 ** (round(np.log2(10 * fs)) + psd_padding))
    window_in_samples = min(int(10 * fs), len(displacement))
    freqs_psd, periodogram = welch(displacement, fs=fs, window='hamming', nperseg=window_in_samples,
                                   noverlap=int(overlap * window_in_samples), nfft=nfft)
    power = 10 * np.log10(periodogram)

    return {'i': i, 'q': q, 'time_vec': time_vec, 'displacement': displacement, 'power': power, 'amplitude': amplitude,
            'freqs': freqs, 'freqs_nlf': freqs_nlf, 'freqs_psd': freqs_psd, 'reference_data': reference_data,
            'raw': data.get('raw'), 'readings': readings, 'channel_select_params': data, 'gap': gap,
            'bin_nums': bin_nums}


def process_single_setup(dont_track, cur_save_path, bp_band, padding, psd_padding, overlap, win_size,
                         update_offset_dynamically, setup_id=None, tlog_path=None, json_path=None, dont_overwrite=None,
                         nlf=False, base_label=''):
    args = get_args()
    bin_nums = None
    if dont_overwrite and images_exists(cur_save_path, dont_track):
        print(f"plot for setup {base_label} exists, skipping...\n")
        return
    data = load_data(setup_id, tlog_path, json_path, args.distance, db)
    gap = data['gap']
    readings = data['data']
    if not isinstance(readings, pd.DataFrame):
        readings = pd.DataFrame(readings)
    if args.time_range is not None:
        readings, bin_nums = cut_signal_by_time_range(readings, bin_nums, gap, base_label)
    original_dist = int(data['dist'])
    config = create_config(setup_id, db)
    if args.select_range or args.variance:
        dr = compute_dr(data, config)
        if args.variance:
            plot_variance(dr, data, base_label, readings)
    if args.select_range:
        bins = bins_select(data, dr)
    else:
        bin_index = fmcw_channel_select(readings.iloc[:int(10 / gap), :], data, config)
        bins = {bin_index: int(data['dist'])}

    bins_dict = {bin_index: process_single_bin(
        bin_index, bins[bin_index], data, readings, gap, dont_track, update_offset_dynamically, setup_id, nlf,
        base_label, bin_nums, bp_band, padding, psd_padding, overlap, win_size) for bin_index in bins.keys()}
    return bins_dict, original_dist


def plot_single_setup(data_dict, cur_save_path, dont_track, band, plots, time_range, setup_id=None, path=None,
                      ppt=None, dont_show=None, intra_breath_path=None):
    label = gen_label(dont_track, setup_id, path, db)

    first_bin = True
    for bin_fmcw in data_dict.keys():
        bin_data = data_dict[bin_fmcw]
        if first_bin:
            plot_raw(data_dict[bin_fmcw].get('raw'), label, cur_save_path, dont_show, ppt)
            plot_bins(bin_data['time_vec'], data_dict[bin_fmcw].get('bin_nums'), label, cur_save_path, dont_show, ppt)
        if 'bokeh' in plots:
            plot_bokeh(data_dict, label, cur_save_path, bin_fmcw, bin_data['displacement'], bin_data['amplitude'],
                       bin_data['time_vec'])
        if 'iq' in plots:
            plot_iq(data_dict, label, cur_save_path, ppt, dont_show, bin_fmcw, time_range, first_bin,
                    len(data_dict.keys()) > 1)
        if 'freqs' in plots:
            plot_freqs(bin_data['freqs'], band, bin_data['reference_data'], f'{label} bin {bin_fmcw}', cur_save_path,
                       ppt, dont_show)
        if bin_data['freqs_nlf'] is not None and 'freqs' in plots:
            if cur_save_path is not None:
                cur_save_path = f'nlf_{cur_save_path}'
            plot_freqs(bin_data['freqs_nlf'], band, bin_data['reference_data'], label + f' nlf bin {bin_fmcw}',
                       cur_save_path, ppt, dont_show)
        if 'psd' in plots:
            plot_psd(bin_data['freqs_psd'], bin_data['power'], bin_data['reference_data'], label, cur_save_path, ppt,
                     dont_show, bin_fmcw)
        if 'iq_vs_time' in plots:
            plot_iq_vs_time(bin_data['time_vec'], bin_data['i'], bin_data['q'],
                            label, cur_save_path, ppt, dont_show, bin_fmcw)
        if 'displacement' in plots:
            plot_displacement(bin_data['time_vec'], bin_data['displacement'], label, cur_save_path, ppt, dont_show,
                              first_bin, setup_id, bin_fmcw, intra_breath_path)
        if 'displacement_lpf' in plots:
            lowcut = back_chair_config['rr']['prep']['low']
            highcut = back_chair_config['rr']['prep']['high']
            phase = Filter().butter_bandpass_filter(TimeArray(bin_data['displacement'], 0.002), lowcut, highcut)
            plot_displacement(bin_data['time_vec'], phase, label + ' Respiration preprocessing', cur_save_path, ppt,
                              dont_show, first_bin, setup_id, bin_fmcw, intra_breath_path)

        if 'amplitude' in plots:
            plot_amplitude(bin_data['time_vec'], bin_data['amplitude'],
                           label, cur_save_path, ppt, dont_show, bin_fmcw)
        first_bin = False


if __name__ == '__main__':
    win_size = 20  # RR
    band = (40 / 60, 250 / 60)
    bp_band = (4 / 60, 300 / 60)
    padding = 2
    psd_padding = 4
    overlap = .9  # 90 percent, for freqs
    args = get_args()
    try:
        db = DB()
    except AttributeError:
        if args.setups is None:
            db = None
        else:
            raise ConnectionError('could not connect to the db')
    change_dir_to_pythondev()
    assert not (args.ppt is not None and args.save_path is None), 'If ppt is inserted, save_path must also be\r\n'
    if args.nlf:
        args.dont_track = True
    if bool(args.setups) + bool(args.tlog_paths) + bool(args.json_paths) == 0:
        master = tk.Tk()
        master.wm_attributes('-topmost', 1)
        tk.Label(master, text="Setup id").grid(row=0)
        e1 = tk.Entry(master)
        e1.grid(row=0, column=1)
        tk.Button(master, text='Go', command=master.quit).grid(row=3, column=0, sticky=tk.W, pady=4)
        tk.mainloop()
        s = e1.get()
        master.destroy()
        if s.isdigit() and len(s) > 0:
            args.setups = [s]
        else:
            db_path = '/Neteera/DATA/' if platform.system() == 'Linux' else 'S:\\'
            ymd = datetime.datetime.now()
            root = tk.Tk()
            root.withdraw()
            root.wm_attributes('-topmost', 1)
            args.tlog_paths = \
                [filedialog.askopenfilename(initialdir=os.path.join(db_path, str(ymd.year), str(ymd.month),
                                                                       str(ymd.day)),
                                               filetypes=(("ttlog files", "*.ttlog"),
                                                          ("tlog files", "*.tlog")))]
    if args.save_path is not None and (not os.path.isdir(args.save_path)):
        print('{} does not exists, creating...'.format(args.save_path))
        os.mkdir(args.save_path)
    if args.ppt is not None:
        ppt = Presentation()
    else:
        ppt = None
    if args.setups:
        db.update_mysql_db(args.setups[0])
        if args.back_twin and not args.variance:
            setu_ids = [find_back_setup_same_session(x, db) for x in args.setups
                        if find_back_setup_same_session(x, db) is not None]
        else:
            setu_ids = []
        setu_ids = sorted(np.unique(setu_ids + args.setups))
        for setu_id in setu_ids:
            cur_save_path = gen_save_path(args.save_path, setu_id)
            data_dict, original_dist = process_single_setup(
                args.dont_track, cur_save_path, bp_band, padding, psd_padding, overlap, win_size=win_size,
                update_offset_dynamically=args.update_offset_dynamically,
                setup_id=setu_id, dont_overwrite=args.dont_overwrite, nlf=args.nlf)
            plot_single_setup(data_dict, cur_save_path, args.dont_track, band=band,  setup_id=setu_id,
                              ppt=ppt, dont_show=args.dont_show, plots=args.plots,
                              time_range=args.time_range, intra_breath_path=args.intra_breath_path)
    elif args.tlog_paths:
        for tlog in sorted(args.tlog_paths):
            cur_save_path = gen_save_path(args.save_path, path=tlog)
            data_dict, original_dist = process_single_setup(
                args.dont_track, cur_save_path, bp_band, padding, psd_padding, overlap,
                update_offset_dynamically=args.update_offset_dynamically, win_size=win_size, tlog_path=tlog,
                dont_overwrite=args.dont_overwrite, nlf=args.nlf, base_label=tlog)
            plot_single_setup(data_dict, cur_save_path, args.dont_track, band=band, path=tlog, ppt=ppt,
                              dont_show=args.dont_show, time_range=args.time_range, plots=args.plots)
    elif args.json_paths:
        for json_path in sorted(args.json_paths):
            cur_save_path = gen_save_path(args.save_path, path=json_path)
            data_dict, original_dist = process_single_setup(
                args.dont_track, cur_save_path, bp_band, padding, psd_padding, overlap,
                update_offset_dynamically=args.update_offset_dynamically, win_size=win_size, json_path=json_path,
                dont_overwrite=args.dont_overwrite, nlf=args.nlf, base_label=json_path)
            plot_single_setup(data_dict, cur_save_path, args.dont_track, band=band, path=json_path, ppt=ppt,
                              dont_show=args.dont_show, time_range=args.time_range, plots=args.plots)

    if args.ppt is not None:
        ppt.save(os.path.join(args.save_path, args.ppt) + '.pptx')
    if args.dont_show:
        plt.close()
        plt.clf()
    else:
        plt.show(block=True)
