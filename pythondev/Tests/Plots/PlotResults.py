# Copyright (c) 2020 Neteera Technologies Ltd. - Confidential

import os
import sys

from Tests.Constants import UNITS
from Tests.Utils.PathUtils import windows_dir_to_linux_dir

sys.path.append(os.path.dirname(os.path.dirname(os.path.dirname(__file__))))  # noqa

import Configurations as Config

from Tests.Utils.DBUtils import calculate_delay
from Tests.Utils.LoadingAPI import load_reference, load_pred_high_qual, get_list_of_setups_in_folder_from_vs
from Tests.Evaluation.EvaluationUtils import under_percent, under_percent_or_thresh, under_thresh
from Tests.Constants import STAT_CLASSES, UNDER_DICT
from Tests.Utils.TestsUtils import *
from Tests.vsms_db_api import DB

import argparse
import numpy as np
import scipy.signal as sp
import matplotlib
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import re
import copy
import warnings

with warnings.catch_warnings():
    warnings.simplefilter("ignore")

VS_SUPPORTED = ['hr', 'rr', 'bbi', 'stat', 'ie']

Y_MIN = {'hr': 39, 'rr': 1, 'stat': -1, 'bbi': 1, 'ie': -1, 'spo2': 80}
Y_MAX = {'hr': 160, 'rr': 50, 'stat': 4, 'bbi': 2000, 'ie': 4, 'spo2': 101}
Y_LIM = {k: (v, Y_MAX[k]) for k, v in Y_MIN.items()}


def get_args() -> argparse.Namespace:
    """ Parse arguments"""
    parser = argparse.ArgumentParser(description='Plot algorithm output versus ground truth')
    parser.add_argument('-result_dir', '-folder_list', metavar='load_path', type=str, required=True,
                        help='Path to load files from')
    parser.add_argument('-save_fig_path', metavar='Location', type=str, required=False,
                        help='location of output figures')
    parser.add_argument('--diff', action='store_true', help='Plot the pred-ref differences')
    parser.add_argument('--overwrite', action='store_true', help='Overwrite existing setups')
    parser.add_argument('-ppt_fname', type=str, required=False, help='Powerpoint filename to generate ppt '
                                                                     'presentation from plots')
    parser.add_argument('-vital_sign', '-compute', nargs='+', type=str, choices=VS_SUPPORTED, default=VS_SUPPORTED,
                        help='plot the following vital signs (default is all)\n')
    parser.add_argument('-match_list_type', metavar='match_lists', type=str, required=False,
                        choices=['NO', 'MSE', 'TS', 'no', 'ts', 'mae'], default='ts',
                        help='what type of match list to use')
    parser.add_argument('--silent', action='store_true', help='Display only warnings and errors')
    parser.add_argument('--t0', action='store_true', help='start each setup from its t0')
    parser.add_argument('-product', type=str, required=False, default='health', help='automotive or health')
    parser.add_argument('--force', action='store_true', help='Process invalid reference setups')
    parser.add_argument('-setups', '-session_ids', metavar='ids', nargs='+', type=int,
                        help='Setup IDs in DB to collect the online results', required=False)
    return parser.parse_args()


def performance(prediction: np.ndarray, ground_truth: np.ndarray, vitalsgn: str,
                perc_or_thresh=None, percent_value=None, thresh_value=None):
    """ Compute performance metric """
    if not len(prediction) == len(ground_truth):
        ground_truth = ground_truth[:min(len(prediction), len(ground_truth))]
        prediction = prediction[:min(len(prediction), len(ground_truth))]
    ab_differences = [(prediction[i], ground_truth[i]) for i, _ in enumerate(prediction)] if \
        set(np.unique(ground_truth)).intersection(np.unique(prediction)) in [{0, 1}, {0}, {1}] else np.abs(
        np.array(prediction) - np.array(ground_truth))
    if vitalsgn == 'hr':
        if perc_or_thresh is None:
            under = under_percent(diffs=ab_differences, ground_truth=ground_truth, percent=10)
        elif perc_or_thresh == 'under_thresh':
            under = under_thresh(diffs=ab_differences, thresh=thresh_value)
        else:
            under = under_percent(diffs=ab_differences, ground_truth=ground_truth, percent=percent_value)
    elif vitalsgn == 'rr':
        if perc_or_thresh is None:
            under = under_thresh(diffs=ab_differences, thresh=4)
        elif perc_or_thresh == 'under_thresh':
            under = under_thresh(diffs=ab_differences, thresh=thresh_value)
        elif perc_or_thresh == 'under_thresh_and_percent':
            under = under_percent_or_thresh(diffs=ab_differences, ground_truth=ground_truth,
                                            thresh=thresh_value, percent=percent_value)
        else:
            under = under_percent(diffs=ab_differences, ground_truth=ground_truth, percent=percent_value)
    else:
        raise ValueError('Invalid vital sign {}'.format(vitalsgn))
    return 100 * np.mean(under) if len(under) else None


def add_margin(vs, gt):
    if vs in ['hr', 'rr']:
        under = UNDER_DICT[vs][0]
        per = under.get('per')
        thresh = under.get('thresh')
        if per is not None:
            if thresh is not None:
                upper = np.maximum(gt + thresh, gt * (1 + per / 100))
                lower = np.minimum(gt - thresh, gt * (1 - per / 100))
            else:
                upper = gt * (1 + per / 100)
                lower = gt * (1 - per / 100)
        else:
            upper = gt + thresh
            lower = gt - thresh
        plt.fill_between(range(len(gt)), upper, lower, alpha=0.5, label='REF margin', color='wheat')


def find_relevant_note(setup_num, db):
    sn_last_digits = db.setup_sn(setup_num)[0][-2:]
    note = db.setup_note(setup_num).replace('\t', '').replace('\n', ', ')
    if 'SN' in note and sn_last_digits in note:
        first_sn = note.find('SN')
        last_sn = note.rfind('SN')
        sn_index = note.find(sn_last_digits)
        if sn_index < last_sn:
            note = note[:last_sn - 1]
            if note.endswith(','):
                note = note[:-1]
        else:
            note = note[:first_sn] + note[last_sn:]
    return note


def add_title_and_information(txt, vs, idw, y_min, y_max, time_until_prediction, match_list_type, db):
    plt.ylim(bottom=y_min, top=y_max)

    plt.title(f'setup: {idw}, Target: {db.setup_target(idw)}, Mount: {db.setup_mount(idw)},'
              f'Posture:{db.setup_posture(idw)}\n{find_relevant_note(idw, db)}')

    distance = db.setup_distance(setup=idw)
    left, right = plt.xlim()
    plt.text(left + 0.01 * (right - left), y_max * 0.86,
             f'{txt}\n'
             f'Distance {distance}[mm]\n'
             f'Delay between reference to Neteera: {time_until_prediction} seconds\n'
             f'Match list type: {match_list_type}', fontsize=10)
    plt.text(left + 0.50 * (right - left), y_max * 0.94, vs, fontsize='xx-large', fontweight='bold')

    plt.grid(True)
    y_label = {'hr': 'HR [bpm]', 'rr': 'RR [bpm]', 'bbi': 'BBI [ms]', 'stat': 'Stat', 'ie': 'Ratio [-]'}
    plt.ylabel(y_label[vs])


def plotter(prediction: np.ndarray,
            gtruth: np.ndarray,
            reliability: np.ndarray,
            vital_sgn: str,
            argss,
            idw: int,
            time_until_prediction: int,
            db,
            base_txt=''):
    """ Plot results against ground-truth    """
    plt.xlabel('Time [s]')

    ax = plt.gca()
    n_ticks = 10
    if vital_sgn == 'bbi':
        jumps = max(1, int(prediction[-1] / n_ticks / 1000))
        ax.set_xticks(np.arange(0, max(prediction[-1], gtruth[-1]) / 1000, jumps))
    if vital_sgn == 'stat':
        plt.yticks(range(len(STAT_CLASSES)), STAT_CLASSES)

    if argss.diff and prediction is not None and gtruth is not None:
        if not len(prediction) == len(gtruth):
            gtruth = gtruth[:min(len(prediction), len(gtruth))]
            prediction = prediction[:min(len(prediction), len(gtruth))]
        plt.plot(prediction - gtruth, label='NETEERA-REF')
        margin = 20
        y_min = min(prediction - gtruth) - margin
        y_max = max(prediction - gtruth) + margin
    else:
        y_min = Y_MIN[vital_sgn]
        y_max = Y_MAX[vital_sgn]
        if gtruth is not None:
            if vital_sgn == 'bbi':
                plt.plot(gtruth[1:] / 1000, np.diff(gtruth), label='REF', color='orange')
            else:
                plt.plot(gtruth, label='REF', color='orange')
        if prediction is not None:
            if vital_sgn == 'bbi':
                plt.plot(prediction[1:] / 1000, np.diff(prediction), label='NETEERA')
            else:
                x_pred = np.arange(len(prediction)) + time_until_prediction
                plt.plot(x_pred, prediction, label='NETEERA', color='royalblue', linewidth=0.65)
            if reliability is not None and len(reliability):
                y_vals = (np.array([r if r > 0 else np.nan for r in reliability]) * prediction[:len(reliability)])
                plt.plot(x_pred, y_vals, color='royalblue', linewidth=2)
        add_margin(vital_sgn, gtruth)

    add_title_and_information(base_txt, vital_sgn, idw, y_min, y_max, time_until_prediction, argss.match_list_type, db)
    if argss.t0:
        t0_values = db.setup_spot(idw)
        if 't0' in t0_values:
            number_of_last_seconds = t0_values['t0']
            plt.axhspan(ymin=0, ymax=y_max, xmin=0, xmax=number_of_last_seconds / len(gtruth),
                        facecolor='red', alpha=0.3)
            plt.text(0, y_max * 0.73, '      t0 = ' + str(number_of_last_seconds), fontsize=10)
    legend = plt.legend(loc='upper right')
    for line in legend.get_lines():
        line.set_linewidth(2.0)

    save_fig_file = f'{vital_sgn}_{idw}_diff.png' if argss.diff else f'{vital_sgn}_{idw}.png'
    if argss.save_fig_path is not None:
        plt.savefig(os.path.join(argss.save_fig_path, save_fig_file))
    else:
        plt.show()
        if argss.ppt:
            slide = argss.ppt.slides.add_slide(argss.ppt.slide_layouts[6])
            image_stream = BytesIO()
            plt.savefig(image_stream)
            slide.shapes.add_picture(image_stream, Inches(2.5), Inches(1.5), height=Inches(3.1))
    plt.close()
    plt.clf()


def handle_tuple(vec):
    if type(vec[0]) is tuple:
        return [p[0] for p in vec]
    else:
        return list(vec)


def finddelay(x, y):
    v = np.zeros(int(round(x[-1]) + 1))
    w = np.zeros(int(round(y[-1]) + 1))
    for p in x:
        v[int(round(p))] = 1
    for g in y:
        w[int(round(g))] = 1

    q = sp.filtfilt(np.ones(5) / 5, [1], v)
    r = sp.filtfilt(np.ones(5) / 5, [1], w)
    s = sp.correlate(r, q)
    return q.shape[0] - np.argmax(s)


def load_gt_pred_reliability(vital_sgn: str, idx: int, argss, db: DB, file: str):
    """ Loads the ground truth data and delays it if necessary"""

    loaded = load_pred_high_qual(os.path.dirname(file), idx, vital_sgn)
    pred = loaded['pred']
    reliability = loaded.get('high_quality')
    is_spot = os.path.basename(os.path.dirname(file)) == 'spot'
    if len(pred) == 0:
        return
    gt = load_reference(idx, str(vital_sgn), db, argss.force)
    if gt is None or len(gt) == 0 or len(pred) == 0:
        gt = np.nan * np.ones(len(pred))
        delay = 0
    else:
        delay = calculate_delay(idx, vital_sgn, db)
    if is_spot and vital_sgn != 'bbi':
        seg_from_end = Config.spot_config['setup']['maximum_window_from_end']
        duration = db.setup_duration(idx)
        if delay is not None and delay > 0:
            gt = gt[delay:]
        gt = gt[int(duration - seg_from_end): int(duration)]
        return {'prediction': pred * len(gt), 'ground_truth': gt, 'reliability': [], 'time_until_prediction': 0}

    if vital_sgn == 'bbi':
        if delay is not None:
            pred += delay * 1000
        gt += finddelay(pred, gt)
        return {'prediction': pred, 'ground_truth': gt, 'reliability': [], 'time_until_prediction': 0}

    if argss.match_list_type.lower() in ['no', 'none', '0'] or isinstance(gt, pd.Series):
        if isinstance(gt, pd.Series):
            gt = gt.to_numpy()
        delay = 0
    if vital_sgn in ['hr', 'rr'] and argss.match_list_type in ['mae', 'mse']:
        _, __, shift = match_lists(pred, gt)
    elif argss.match_list_type == 'ts':
        shift = delay
    else:
        shift = 0

    return {'prediction': pred, 'ground_truth': gt, 'reliability': reliability,
            'time_until_prediction': shift}


def load_and_plot(setup, vital_signn, argss, db, file):
    print(f'plotting file: {file}')
    return_dict = load_gt_pred_reliability(vital_signn, setup, argss, db, file)
    if return_dict is not None:
        plotter(return_dict['prediction'],
                return_dict['ground_truth'],
                return_dict['reliability'],
                vital_signn,
                argss,
                setup,
                return_dict['time_until_prediction'],
                db)


def plot_bland_altman(y_ref, y_pred, vs, save_path, title='', x_axis_ref=False):
    """
    :param y_ref:   reference true values
    :param y_pred:  prediction values
    :param vs:  vital sign
    :param save_path:   path to save the plot
    :param title:   title of the plot
    :param x_axis_ref:  Detrmine if the x_axis is y_ref (default is mean(y_ref, y_pred)
    """
    predictions, gtruths = np.array(y_pred), np.array(y_ref)
    mean = np.nanmean([predictions, gtruths], axis=0)
    x_axis = gtruths if x_axis_ref else mean
    x_label = 'ref' if x_axis_ref else 'mean'
    diff = predictions - gtruths
    md = np.nanmean(diff)
    sd = np.nanstd(diff, axis=0)
    plt.scatter(x_axis + np.random.normal(0, 0.2, len(mean)), diff, facecolors='none', edgecolors='b')
    plt.xlabel(f'{x_label} [{UNITS[vs]}]')
    plt.ylabel(f'pred - reference [{UNITS[vs]}]')
    plt.axhline(md, color='gray', linestyle='--')
    plt.axhline(md + 1.96 * sd, color='gray', linestyle='--')
    plt.axhline(md - 1.96 * sd, color='gray', linestyle='--')
    plt.title(f'{title} Bland Altman plot vital sign {vs} number of points: {len(mean)}')
    plt.savefig(save_path)
    plt.close()
    plt.clf()


def load_and_plot_bland_altman(setup_nums, vs, save_path, args, db):
    preds, gts = [], []
    for setup_num in setup_nums:
        setup_dict = load_gt_pred_reliability(
            vs, setup_num, args, db, os.path.join(args.result_dir, f'{setup_num}_{vs}_spot.data'))
        preds.append(np.median(setup_dict['prediction']))
        gts.append(np.median(setup_dict['ground_truth']))
    plot_bland_altman(gts, preds, vs, save_path)


def plot_single_vs_single_folder(args, vs, db):
    """plot a single vital sign for a given folder, return whether result file were found"""

    pred_files = get_list_of_setups_in_folder_from_vs(args.result_dir, vs)
    if args.setups is not None:
        pred_files = {setup: file for setup, file in pred_files.items() if setup in args.setups}
    if len(pred_files) == 0:
        return False
    spot_setup_nums = [re.findall(fr"[\d]+_{vs}_spot", file) for file in pred_files.values()]
    spot_setup_nums = [int(re.findall(r'\d+', x[0])[0]) for x in spot_setup_nums if len(x)]
    if len(spot_setup_nums) > 0 and vs in ['hr', 'rr']:
        load_and_plot_bland_altman(spot_setup_nums, vs, os.path.join(args.save_fig_path, f'bland_altman_{vs}.png'), args, db)
    for setup, file in pred_files.items():
        load_and_plot(setup, vs, args, db, os.path.join(args.result_dir, file))
    return True


def load_and_plot_single_result_file(argss, file, db):
    idxx = int(re.findall(r"[\d]+", os.path.basename(file))[0])  # extracts setup id from result file name
    if file.endswith('.csv'):
        for vs in argss.vital_sign:
            load_and_plot(idxx, vs, argss, db, file)
    else:
        file_name = file.split('/')[-1]
        file_name = os.path.basename(file_name).split('.')[0]
        vs = file_name if file_name.isnumeric() else file_name.split('_')[1]
        if vs in argss.vital_sign:
            load_and_plot(idxx, vs, argss, db, file)


def plot_single_folder(args, db):
    if args.ppt_fname is not None:
        ppt = Presentation()
        title_slide_layout = ppt.slide_layouts[6]
        args.ppt = ppt
    else:
        args.ppt = None
    assert os.path.exists(args.result_dir)
    if os.path.isdir(args.result_dir):
        if not is_debug_mode() or 'PlotResults' not in __file__:
            matplotlib.use('Agg')
        if not args.save_fig_path:
            args.save_fig_path = args.result_dir
        is_found_files = 0
        for vs in args.vital_sign:
            is_found_files += plot_single_vs_single_folder(args, vs, db)
        if not is_found_files:
            raise FileNotFoundError(f'could not find pred files in {args.result_dir}')
    else:
        load_and_plot_single_result_file(args, args.result_dir, db)


def main_plot_results(arguments, db):
    arguments.match_list_type = arguments.match_list_type.lower().replace('mse', 'mae')
    if platform.system() == 'Linux':
        arguments.result_dir = windows_dir_to_linux_dir(arguments.result_dir)
    try:
        collect_result(copy.deepcopy(arguments.setups), os.path.dirname(arguments.result_dir))
        plot_single_folder(arguments, db)
    except FileNotFoundError:
        sub_folders = [x[0] for x in os.walk(arguments.result_dir) if os.path.basename(x[0]) in ['dynamic', 'spot']]
        for sub_dir in sub_folders:
            print('going to sub-folder ', sub_dir)
            arguments.result_dir = sub_dir
            arguments.save_fig_path = sub_dir
            plot_single_folder(arguments, db)


if __name__ == "__main__":
    main_plot_results(get_args(), DB())
