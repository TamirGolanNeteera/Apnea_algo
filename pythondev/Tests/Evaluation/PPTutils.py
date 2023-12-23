
from Tests.Utils.TestsUtils import get_setups_list
from Tests.Utils.StringUtils import join
from Tests.Evaluation.pd2ppt import df_to_table

import os
from pptx import Presentation
from pptx.util import Cm, Pt
import numpy as np
import pandas as pd
import itertools

IMG_LOCATIONS_VERTICAL = np.array([0.95, 9.3])
IMG_LOCATIONS_HORIZONTAL = [5.2, 15]

TEXT_LOCATIONS_VERTICAL = IMG_LOCATIONS_VERTICAL - 0.7
TEXT_LOCATIONS_HORIZONTAL = [13.5, 15.3]
HEIGHT = 7.9
WIDTH = 5.5

IMG_LOCATIONS_VERTICAL_MANY = np.array([0, 6.2, 12.6])
IMG_LOCATIONS_HORIZONTAL_MANY = [7.7, 16.5]

TEXT_LOCATIONS_VERTICAL_MANY = IMG_LOCATIONS_VERTICAL_MANY - 0.7
TEXT_LOCATIONS_HORIZONTAL_MANY = [14.4, 16.9]
HEIGHT_MANY = 6.1
WIDTH_MANY = 8

ERROR_DIFFS_LOC = [5.5, 15]
ERROR_DIFFS_LOC_MANY = [0, 19]


def get_locations(number_of_version) -> (list, list, float, float):
    if number_of_version <= 4:
        img_horizontal = IMG_LOCATIONS_HORIZONTAL
        img_vertical = IMG_LOCATIONS_VERTICAL
        text_horizontal = TEXT_LOCATIONS_HORIZONTAL
        text_vertical = TEXT_LOCATIONS_VERTICAL
        height = HEIGHT
        width = WIDTH
    else:
        img_horizontal = IMG_LOCATIONS_HORIZONTAL_MANY
        img_vertical = IMG_LOCATIONS_VERTICAL_MANY
        text_horizontal = TEXT_LOCATIONS_HORIZONTAL_MANY
        text_vertical = TEXT_LOCATIONS_VERTICAL_MANY
        height = HEIGHT_MANY
        width = WIDTH_MANY
    diff_locations = ERROR_DIFFS_LOC if number_of_version < 3 else ERROR_DIFFS_LOC_MANY
    return (list(itertools.product(img_vertical, img_horizontal)),
            list(itertools.product(text_vertical, text_horizontal)),
            height,
            width,
            diff_locations)


def add_text_box(s, text, left, top, font=None):
    text_box = s.shapes.add_textbox(Cm(left), Cm(top), Cm(1), Cm(1))
    text_box.text_frame.text = str(text)
    if font is not None:
        for paragraph in text_box.text_frame.paragraphs:
            paragraph.font.size = Pt(font)


def handle_db_view_entry(x):
    if isinstance(x, list):
        if len(set(x)) == 1:
            return list(set(x))[0]
        return join(x, '\n')
    return x


def add_single_slide(prs, view, idx, ver_names, fig_paths, error_diffs=dict()):
    img_locations, text_locations, img_size, width_table, error_diff_loc = get_locations(len(ver_names))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i, ver in enumerate(fig_paths):
        slide.shapes.add_picture(fig_paths[i], Cm(img_locations[i][1]), Cm(img_locations[i][0]), height=Cm(img_size))
    for i, ver in enumerate(ver_names):
        horizontal_loc = text_locations[i][1]
        if i % 2 == 0:
            horizontal_loc -= 0.19 * len(str(ver))
        add_text_box(slide, ver, horizontal_loc, text_locations[i][0], 14)  # version name

    if view is not None and idx in view.index:
        series = view.loc[int(idx), ['Setup Type', 'Video', 'session', 'time', 'duration', 'subject', 'gender', 'path',
                                     'location', 'target', 'posture', 'state', 'distance', 'sn', 'notes', 'validity']]
        df = pd.DataFrame(series.apply(handle_db_view_entry))
        df_to_table(slide, df.dropna().reset_index(), left=0, top=0, width=width_table)
    add_text_box(slide, error_diffs.get(int(idx), ''), *error_diff_loc, 14)  # error diff
    add_text_box(slide, idx, 23.5, 18, 14)  # page setup id


def create_compare_fig_ppt(dirs, versions, output_name, vs, setup_ids=None, db_view=None, metadata={}, min_length=0):
    if setup_ids is None:
        setups_lists = [set(get_setups_list(folder, vs)) for folder in dirs]
        setup_ids = set.intersection(*setups_lists)
        if not setup_ids:
            return
    if len(setup_ids) < min_length:
        return
    presentation = Presentation()
    for setup in sorted(setup_ids):
        fig_full_paths = [os.path.join(folder, f'{vs}_{setup}.png') for folder in dirs]
        fig_full_paths = [x for x in fig_full_paths if os.path.isfile(x)]
        add_single_slide(presentation, db_view, setup, versions, fig_full_paths, metadata)
    if not output_name.endswith('.pptx'):
        output_name = os.path.join(os.path.dirname(dirs[0]), f'{output_name}_{vs}.pptx')
    presentation.save(output_name)


def create_compare_ppt_between_radars(folder, feature, output_name, vs, db_view, metadata={}, multi_vs=True):
    if db_view is None:
        print('can\'t make compare between radars')
        return
    presentation = Presentation()
    for session, row in db_view.T.items():
        row_df = pd.DataFrame(row[['setup', feature, 'sn']].to_dict()).sort_values(feature, ascending=False)
        feature_labels = [f'{setup}; {feature} = {value}  SN: {sn[-3:]}' for setup, value, sn in row_df.values]
        fig_full_paths = [os.path.join(folder, f'{vs}_{setup}.png') for setup in row_df['setup'].values]
        if multi_vs and vs == 'hr' and len(feature_labels) <= 2:
            feature_labels += feature_labels
            fig_full_paths += [x.replace('hr_', 'rr_') for x in fig_full_paths]
        fig_full_paths = [x for x in fig_full_paths if os.path.isfile(x)]

        if len(fig_full_paths):
            add_single_slide(presentation, db_view, session, feature_labels, fig_full_paths, metadata)
    if len(presentation.slides):
        presentation.save(os.path.join(os.path.dirname(folder), f'{vs}_{feature}_{output_name}.pptx'))
    else:
        print('empty pptx, not saving')
