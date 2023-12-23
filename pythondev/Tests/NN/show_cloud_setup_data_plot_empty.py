# Copyright (c) 2021 Neteera Technologies Ltd. - Confidential
import os
import sys
sys.path.insert(1, os.getcwd())
# from Tests.Plots.PlotRawDataRadarCPX import*
import argparse
import fnmatch
from os import listdir
from sklearn import preprocessing
import pandas as pd
from Tests.Utils.LoadingAPI import load_reference
from Tests.vsms_db_api import *
from pptx import Presentation
from pptx.util import Inches, Cm
from io import BytesIO
from Tests.NN.create_apnea_count_AHI_data import get_signal_time_index, MB_HQ, delays, getSetupDataLocalDB, getSetupRespirationCloudDBDebug, getSetupRespirationCloudDBDebugWithTimestamp, MB_HQ, count_apneas_in_chunk, getSetupRespirationCloudDBDebugDecimate, getSetupRespirationLocalDBDebug, getSetupRespirationCloudDB, compute_respiration, compute_phase
from Tests.NN.Chest.create_apnea_count_AHI_data_regression_MB_chest import apnea_class
import matplotlib.pyplot as plt
import glob
import scipy.signal as sp
from scipy.signal import spectrogram
from Tests.NN.reiniers_create_data import spectrogram_features
from scipy.signal import butter, filtfilt, lfilter, firwin
import datetime
import matplotlib.pyplot as plt
import matplotlib.dates as mdates

db = DB()
home = '/Neteera/Work/homes/dana.shavit/Research/analysis/'



def find_lowest_threshold_index(A, N):
    # Initialize variables to store the minimum threshold and its index
    min_threshold = float('inf')  # Initialize with positive infinity
    min_index = -1

    # Iterate through the thresholds and find the lowest one greater than N
    for i, threshold in enumerate(A):
        if threshold > N and threshold < min_threshold:
            min_threshold = threshold
            min_index = i

    return min_index

ss_dict = {'N1':1, 'N2':1, 'N3':3, 'R':4, 'W':0, 'invalid':-1}

bin_dict = {}
for i in range(-1,15):
    bin_dict[i] = []

amp_keys = [250, 500, 750, 1000, 2000, 3000, 4000, 5000, 6000, 7000, 8000, 9000, 10000]
amp_dict = {}
for k in amp_keys:
    amp_dict[k] = []


def get_sleep_stage_labels(setup:int):
    ss_ref = []
    ss_ref = load_reference(setup, 'sleep_stages', db)

    if ss_ref is None:
        return []

    if isinstance(ss_ref, pd.core.series.Series):
        ss_ref = ss_ref.to_numpy()

    #ss_ref_class = np.zeros_like(ss_ref)
    ss_ref_class = [ss_dict[ss] for ss in ss_ref]

    return ss_ref_class





def rollavg_convolve_edges(a,n):
    'scipy.convolve, edge handling'
    assert n%2==1
    return sp.convolve(a,np.ones(n,dtype='float'), 'same')/sp.convolve(np.ones(len(a)),np.ones(n), 'same')


setups = ['all']
#base_path = '/Neteera/Algo/david.grossman/Results/2023/08/6/original_thresh_0.4/stat_06_08_2023_es_office_benchmark/'
#base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.5.11.2/angle_test_0708_no_ls_selected_sessions/'
#base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.5.11.2/no_ls_szmc_0708_0.28_motion_selected_sessions/'
base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.5.11.2/no_ls_herzog_0908_selected_sessions/'
#base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.5.11.2/aaa_merge_1308_cloud_selected_sessions/'
#base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.5.11.2/aaa_merged_version_1308_usb_selected_sessions/'
show_dupes = False
#base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.5.11.2/cloud_debug_smoothing_selected_sessions/'#angle_tests_130_0808_2/'#210_angle_test_selected_sessions/'

# base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.6.xx/szmc_1hz/'
# base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.6.xx/latest/'
base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.6.xx/AAC_2/'
base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.6.xx/lab_setups_no_binchange5/'
#base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.6.xx/lab_setups/'#_no_binchange/'
#base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.6.xx/lab_setups_111310/'
base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.6.xx/braintree/'
base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.6.xx/aac_3_6_0_lowrr/'#aac_3_6_0/'#aac_with_sig_qual/'
#base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.6.xx/aac_with_sig_qual/'#aac_3_6_0_lowrr/'#aac_3_6_0/'#aac_with_sig_qual/'
#base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.6.xx/szmc_with_sq_selected_sessions/'
#base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.6.xx/szmc_new_2_selected_sessions/'
base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.6.xx/aac_3_6_0_shimon_lowrr_rules2/'#aac_3_6_0_new_lowrr_rules/'#aac_rr_conditions/'
base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.6.xx/aac_3_6_0_high_rr_with_rules_and_mqual/'#aac_3_6_0_new_lowrr_rules/'#aac_rr_conditions/'
base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.6.xx/aac_3_6_0_low_signal_2500/'
base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.6.xx/bp/'
#base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.6.xx/mb_1211/'#aac_3_6_0_low_signla/'
outfolder=base_path

# Location	Location #/group	                                                Raleigh Room 1	Raleigh Room 3	Sumter
# Room 1	Sumter Room 2	Sumter Room 3 Straight above the patient at 80cm.	                        1
# 213023020232	213023020238	213023020234	213023020236	213023020231 Parallel to group 1 on opposite side of
# bed at 80cm.	        2	            213023020240	213023020248	213023020254	213023020250	213023020251
# Straight above the patient at 1.5m.	               _bin         3	            213023020270	213023020268
# 213023020256	213023020269	213023020259 Off the side of the bed, on the side that the psg is located	4
# 213023020278	213023020279	213023020273	213023020271	213023020274

radar_dict = np.array([[213023020232,	213023020238,	213023020234,	213023020236,	213023020231],
                       [213023020240,	213023020248,	213023020254,	213023020250,	213023020251],
                       [213023020270,	213023020268,	213023020256,	213023020269,	213023020259],
                       [213023020278,	213023020279,	213023020273,	213023020271,	213023020274]])
rooms = ["Raleigh Room 1",	"Raleigh Room 3",	"Sumter Room 1",	"Sumter Room 2",	"Sumter Room 3"]
locations = [1,2,3,4]
descriptions = ["Straight above the patient at 80cm", "Parallel to group 1 on opposite side of bed at 80cm.", "Straight above the patient at 1.5m", "Off the side of the bed, on the side that the psg is located"]

setups = ['all']
setups = [112135,112136, 112137,112138, 112139, 113140,112141]
#setups = [112047, 112048, 112049, 112050]
shw=True
if setups == ['all']:
    shw=True
    setups = []
    setup_files = fnmatch.filter(os.listdir(os.path.join(base_path,'dynamic')), '*_hr.npy')

    print("setup_files", setup_files)

    for i, fn in enumerate(setup_files):
        sess = int(fn[0:fn.find('_')])
        setups.append(sess)

print(setups)

def high_pass_bin(signal):

    # Apply high-pass filter with cutoff frequency of 20 Hz
    cutoff_freq = 0.1  # Cutoff frequency (Hz)
    fs = 500
    # Design the high-pass filter
    nyquist_freq = 0.5 * fs
    norm_cutoff_freq = cutoff_freq / nyquist_freq
    b, a = butter(3, norm_cutoff_freq, btype='high', analog=False, output='ba')

    # Apply the filter to the signal
    filtered_signal = filtfilt(b, a, signal)

    return filtered_signal
def get_args() -> argparse.Namespace:
    """ Argument parser

    :return: parsed arguments of the types listed within the function
    :rtype: argparse.Namespace
    """
    parser = argparse.ArgumentParser(description='Process some integers.')
    parser.add_argument('-save_path', metavar='Location', type=str, required=True, help='location of saved data')
    parser.add_argument('--scale', action='store_true', help='Scale data to m=0 s=1')
    parser.add_argument('--show', action='store_true', help='Scale data to m=0 s=1')
    parser.add_argument('--overwrite', action='store_true',  required=False, help='Overwrite existing output')
    parser.add_argument('-chunk', metavar='window', type=int, required=True, help='signal size')
    parser.add_argument('-step', metavar='window', type=int, required=True, help='stride for signal creation')
    return parser.parse_args()

cols = {'empty chair': 'red', 'empty':'red',
                    'low_s': 'yellow',
                    'motio': 'blue',
                    'runni': 'green',
                    'full': 'green',
                    'unkno': 'grey',
                    'warmi': 'magenta',
                    'missi': 'black',
                    'low_r':'black'}
cols_class = {0: 'red', 1: 'yellow', 2: 'blue', 3: 'green', 4: 'green', 5: 'grey', 6: 'olivedrab', 7: 'black', 8:'black'}
to_draw = {'empty':True,
                    'low_s': True,
                    'motio': True,
                    'runni': True,
                    'full': True,
                    'unkno': True,
                    'warmi': True,
                    'missi': False,
                    'low_r':False}
to_class = {'empty':0,
                    'low_s': 1,
                    'motio': 2,
                    'runni': 3,
                    'full': 3,
                    'unkno': 4,
                    'warmi': 5,
                    'missi': 6,
                    'low_r':7}

def standardize_key(curr_key):
    s_key = curr_key.lower()
    s_key = s_key.replace(' ', '_')
    s_key = s_key.replace('-', '_')
    s_key = s_key[:min(5, len(s_key))]
    return s_key


def getSegments(data_frame, seg_val, decimate=False):
    #
    countr = 0
    segments = []
    for k, v in data_frame[data_frame['stat'] == seg_val].groupby((data_frame['stat'] != seg_val).cumsum()):
        if decimate:
            segments.append([int(np.floor(v.index[0])), 1+int(np.floor(v.index[-1]))])
        else:
            segments.append([int(v.index[0]), 1+int(v.index[-1])])
        countr += 1
    return segments


if __name__ == '__main__':
    col = ['gray', 'blue', 'green', 'red', 'yellow', 'magenta', 'cyan']

    print(setups)
    #setups = [109870]
    d = []
    ppt = Presentation()

    binchanges_list = []
    harmony_dict = {}
    false_lowrr = {}
    title_slide_layout = ppt.slide_layouts[6]
    sum_10_after = 0
    sum_10_before = 0
    sum_30_after = 0
    sum_30_before = 0

    for i_sess, sess in enumerate(setups):
        db.update_mysql_db(sess)
        # if db.setup_distance(sess)< 450:
        #     continue
        online_10 = 0
        offline_10 = 0
        online_30 = 0
        offline_30 = 0
        online_10_neg = 0
        offline_10_neg = 0
        online_30_neg = 0
        offline_30_neg = 0
        ts = None
        use_ts = False

        # import scipy.signal as sp
        # try:
        #     respiration, fs_new, bins, I, Q = getSetupRespirationCloudDBDebug(sess)
        #     plt.figure()
        #
        #     f, times, Sxx = spectrogram(I.flatten(), fs=fs_new, nperseg=30 * fs_new, noverlap=5 * fs_new)
        #     fimax = np.where(f > 12)[0][0]  ## Max 30 Hz
        #     f = f[2:fimax]
        #     Sxx = Sxx[2:fimax, :]
        #
        #     plt.pcolormesh(times, f, 10 * np.log10(Sxx))  # Convert to dB scale for better visualization
        #     #plt.show()
        #     plt.savefig(outfolder + str(sess) + "_I_spec.png")
        #
        #     plt.close()
        #     plt.figure()
        #     f, times, Sxx = spectrogram(Q.flatten(), fs=fs_new, nperseg=30 * fs_new, noverlap=5 * fs_new)
        #     fimax = np.where(f > 12)[0][0]  ## Max 30 Hz
        #     f = f[2:fimax]
        #     Sxx = Sxx[2:fimax, :]
        #
        #     plt.pcolormesh(times, f, 10 * np.log10(Sxx))  # Convert to dB scale for better visualization
        #     # plt.show()
        #     plt.savefig(outfolder + str(sess) + "_Q_spec.png")
        #     plt.close()
        #     continue
        # except:
        #     continue
        fig, ax = plt.subplots(6, sharex=True, figsize=(20, 10))
        try:
            if sess < 100000:
                respiration, fs_new, bins, I, Q = getSetupDataLocalDB(sess)
                print("local setup", sess, "duration", db.setup_duration(sess))

                # try:
                #     fig1, ax1 = plt.subplots(5, sharex=True, figsize=(14, 7))
                #     # for i in range(5):
                #     #     ax1[i].plot(bins[:,i].real, linewidth=0.25, label='I bin ')
                #     #     ax1[i].plot(bins[:,i].imag, linewidth=0.25, label='Q bin ')
                #     #     #ax1.legend()
                #     # print("saving bins 1")
                #     # plt.savefig(outfolder + str(sess) + "_bins_1.png")
                #     # plt.close(fig1)
                #     # fig2, ax2 = plt.subplots(5, sharex=True, figsize=(14, 7))
                #     # for i in range(5):
                #     #     ax2[i].plot(bins[:, 5+i].real, linewidth=0.25, label='I bin ')
                #     #     ax2[i].plot(bins[:, 5+i].imag, linewidth=0.25, label='Q bin ')
                #     #     #ax2.legend()
                #     # print("saving bins 2")
                #     # plt.savefig(outfolder + str(sess) + "_bins_2.png")
                #     # plt.close(fig2)
                #     # fig3, ax3 = plt.subplots(5, sharex=True, figsize=(14, 7))
                #     # for i in range(5):
                #     #     ax3[i].plot(bins[:, 10+i].real, linewidth=0.25, label='I bin ')
                #     #     ax3[i].plot(bins[:, 10+i].imag, linewidth=0.25, label='Q bin ')
                #     #     #ax3.legend()
                #     # print("saving bins 3")
                #     # plt.savefig(outfolder + str(sess) + "_bins_3.png")
                #     # plt.close(fig3)
                #
                # except:
                #     # plt.close()
                #     continue
                # continue
            if sess > 100000:

                respiration, fs_new, bins, I, Q = getSetupRespirationCloudDBDebug(sess)
                # respiration = np.load(
                #     "/Neteera/Work/homes/tamir.golan/embedded_phase/Tamir_sessions/hr_rr_ra_ie_stat_intra_breath_26_09_2023/accumulated/111577_stat_phase.npy",
                #     allow_pickle=True)
                #respiration = np.concatenate(respiration.to_numpy())
                print(len(respiration), len(I), len(Q))
                print("LENGTH BINS", len(bins))
                binss = np.repeat(bins, fs_new)
                ax[3].plot(respiration, label='resp', linewidth=0.75)
        except:
            print("setup not found, skipping")
            continue
        print(fs_new, "HZ")
        ts =  np.arange(0, db.setup_duration(sess), 1/fs_new)
        print("len(ts)", len(ts))
        print("len(I)", len(I))

        session = db.session_from_setup(sess)
        dist = db.setup_distance(sess)

        try:
            ax[5].plot(binss, label='bins', linewidth=0.75)
        except:
            print("no bins")

        ax[5].set_title(str(sess)+' Selected Bins')

        ax[1].plot(I, label='I online', linewidth=0.25)
        ax[1].plot(Q, label='Q online', linewidth=0.25)

        try:
            p = db.setup_dir(sess)
            raw_dir = os.sep.join([p, 'NES_RAW'])
            res_dir = os.sep.join([p, 'NES_RES'])
        except:
            continue

        plot_sleep_wake = False
        if plot_sleep_wake:
            try:
                print("loading sleep wake data")
                ss = get_sleep_stage_labels(sess)
                ss = np.repeat(ss, fs_new)
                print("done drawing df['stat']")
                ax[4].plot(ss)

            except:
                print("failed to do s/w on", sess)

        hr_online_by_sec = []
        rr_online_by_sec = []
        hr_offline_by_sec = []
        rr_offline_by_sec = []
        try:
            if 'cloud' in db.mysql_db:
                onlyfiles = listdir(res_dir)
                print(res_dir)
                vs = [f for f in onlyfiles if 'VS' in f and 'csv' in f]

                if len(onlyfiles) ==1:
                    res_dir = os.sep.join([res_dir, onlyfiles[0]])
                    onlyfiles = listdir(res_dir)
                    vs = [f for f in onlyfiles if 'results' in f and 'csv' in f]

                vs = vs[0]

                df = pd.read_csv(os.sep.join([res_dir, vs]))
            else:
                p = db.setup_dir(sess)
                res_dir = os.sep.join([p, 'NES_RES'])
                onlyfiles = listdir(res_dir)
                print(onlyfiles)
                if len(onlyfiles)>1:
                    print("!!!")
                res_dir = os.path.join(res_dir, onlyfiles[0])
                onlyfiles = listdir(res_dir)
                vs = [f for f in onlyfiles if 'results' in f and 'csv' in f]
                vs = vs[0]

                df = pd.read_csv(os.sep.join([res_dir, vs]))

            rr = np.repeat(df['rr'].to_numpy(), fs_new)
            hr_online_by_sec = df['hr'].to_numpy()
            rr_online_by_sec = df['rr'].to_numpy()
            setup_length = len(rr)
            rr_hq = df['rr_hq'].to_numpy()
            hr = np.repeat(df['hr'].to_numpy(), fs_new)
            hr_hq = df['hr_hq'].to_numpy()
            print("LENGTH HR HQ", len(hr_hq))

            df['stat'] = df['stat'].fillna('missing')
            #status = df['stat'].to_numpy()
            #print(np.unique(status))
            print(len(df), len(I)/fs_new)

            hr_hq_rate_online = int(np.round(100 * len(hr_hq[hr_hq==1])/len(hr_hq[hr_hq>=0])))
            rr_hq_rate_online = int(np.round(100 * len(rr_hq[rr_hq==1])/len(rr_hq[rr_hq>=0])))

            rr_hq_rate = int(np.round(100 * len(rr_hq[rr_hq==1])/len(rr_hq[rr_hq>=0])))
            hr_val_rate_online = int(np.round(100 * len(hr[hr>0])/len(hr)))
            rr_val_rate_online = int(np.round(100 * len(rr[rr>0])/len(rr)))
            print(hr_hq_rate_online, rr_hq_rate)
            ax[2].set_title('VS')
            ax[2].plot(hr, linewidth=1, alpha=0.3, color='red', label='HR: HQ '+str(hr_hq_rate_online)+'%,  Pos: '+str(hr_val_rate_online)+'%')
            ax[2].plot(rr, linewidth=1, alpha=0.3, color='blue', label='RR: HQ '+str(rr_hq_rate_online)+'%,  Pos: '+str(rr_val_rate_online)+'%')

            for s in df['rr'].to_numpy():

                if s > 5 and s < 10:
                    online_10 += 1
                if s < -5 and s > -10:
                    online_10_neg += 1
                if s > 32:
                    online_30 +=1
                if s < -32:
                    online_30_neg +=1

            print(np.unique(df['stat']))
            empty_online = 0
            stat_empty_online = None
            stat_empty_online = (df['stat'] == 'Empty').to_numpy()
            for uu in np.unique(df['stat']):
                vv = df['stat'] == uu
                stat_vec = np.zeros_like(vv)
                stat_vec[vv == True] = 1

                print("Percentage", uu, 100 * (sum(vv) / len(vv)))
                print("df['stat'] getting segments for", uu)
                es = getSegments(df, uu)

                #print("ONLINE", uu, es)
                if 'mpty' in uu or 'MPTY' in uu:
                    empty_online = len(es)

                print("success...")
                # ax[1].set_title('Device Online Computed Status')
                for e in es:

                    suu = standardize_key(uu)
                    if suu in to_draw.keys() and to_draw[suu]:
                        try:

                            ax[1].axvspan(e[0] * fs_new, e[1] * fs_new, color=cols[suu], alpha=0.3)
                        except:
                            print("plotting", uu, "failed")
                print("done drawing df['stat']")
        except:
            print("no online result data for setup", sess)


        if show_dupes:
            dupes = []
            nsec = int(np.floor(len(I) / fs_new))
            new_len = int(nsec * fs_new)
            I = I[:new_len]
            sig_seconds = I.reshape(nsec, fs_new)

            firstsec = 0
            lastsec = sig_seconds.shape[0]
            for i in range(firstsec, lastsec):
                if i % 1000 == 0:
                    print(i, "/", lastsec)
                pattern1 = sig_seconds[i, 0:10]
                p1 = np.tile(pattern1, [nsec, 1])

                diff = np.sum(np.abs(sig_seconds[:, 0:10] - p1), axis=1)

                zero_idxs = np.where(diff < 0.000001)[0]

                if len(zero_idxs) > 1:
                    # np.delete(zero_idxs, np.where(zero_idxs == i))
                    dupes.append(zero_idxs)

                    for z in zero_idxs:
                        ax[0].axvspan((z * fs_new), (z + 1) * fs_new, color='magenta', alpha=0.2)
                        ax[1].axvspan((z * fs_new), (z + 1) * fs_new, color='magenta', alpha=0.2)
                        ax[2].axvspan((z * fs_new), (z + 1) * fs_new, color='magenta', alpha=0.2)
                        ax[3].axvspan((z * fs_new), (z + 1) * fs_new, color='magenta', alpha=0.2)
            sn = int(db.setup_sn(sess)[0])
            print(dupes)

        binchanges=0
        #empty chair reasons
        nt = [] if not db.setup_note(sess) else db.setup_note(sess)
        plot_empty_data = True
        if plot_empty_data:
            reason = None
            try:
                print("estimated range")
                seleced_bins = np.array(np.load(os.path.join(base_path + 'accumulated', str(sess) + '_bin.npy'),
                                          allow_pickle=True))
                print(seleced_bins)
                seleced_bins = np.array([-1 if b == None else b for b in seleced_bins])
                bd = np.diff(seleced_bins,prepend=0)
                #bd[stat_empty_online == False] = 0
                binchanges=len(bd[bd!=0])
                binchanges_list.append([sess ,binchanges])
                seleced_bins = np.repeat(seleced_bins, fs_new)
                # for u in np.unique(reason):
                #     print(u, 100*len(reason[reason == u])/len(reason))
                #ax[3].plot(seleced_bins, color='magenta', label='range, changes='+str(binchanges), linewidth=0.5)
                # ax[3].set_yticks(range(-1, 5, 1))
            except:
                print('no qual fail data for setup', sess)
            reason = None
            try:
                #print("loading raw rr quality")
                reason = np.array(np.load(os.path.join(base_path + 'accumulated', str(sess) + '_zc.npy'),
                                          allow_pickle=True))
                reason = np.repeat(reason, fs_new)
                ax[4].plot(reason, color='magenta', label='zc %ile', linewidth=0.5)

                print("loading zero crossing data")
                reason = np.array(np.load(os.path.join(base_path + 'accumulated', str(sess) + '_mzc.npy'),
                                          allow_pickle=True))
                reason = np.repeat(reason, fs_new)

                ax[4].plot(reason, color='black', label='zc %ile smooth', linewidth=0.5)
                mean_noise = np.mean(reason)


            except:
                print('no zc data for setup', sess)

            sig_qual=None
            try:
                #print("loading raw rr quality")
                sig_qual = np.array(np.load(os.path.join(base_path + 'accumulated', str(sess) + '_qual.npy'),
                                          allow_pickle=True))
                sig_qual = np.repeat(sig_qual, fs_new)
                ax[3].plot(sig_qual, color='magenta', label='signal quality', linewidth=0.5, alpha = 0.6)

                print("loading zero crossing data")
                sig_qual = np.array(np.load(os.path.join(base_path + 'accumulated', str(sess) + '_mqual.npy'),
                                          allow_pickle=True))
                sig_qual = np.repeat(sig_qual, fs_new)
                ax[3].plot(sig_qual, color='black', label='signal quality', linewidth=0.5, alpha = 0.6)
                #ax[5].plot(sig_qual, color='black', label='signal quality smooth', linewidth=0.5)
                ax[3].axhline(y=0.5, color='magenta', linewidth=0.5)
                #ax[5].legend()
            except:
                print('no signal qualty data for setup', sess)
            reason = None
            try:
                #print("loading raw rr data")
                #reason = np.array(np.load(os.path.join(base_path + 'accumulated', str(sess) + '_rr_raw.npy'),
                #                          allow_pickle=True))
                #reason = np.repeat(reason, fs_new)
                #ax[3].plot(reason, color='magenta', label='raw rr', linewidth=0.5)

                print("loading amp data")
                amp_sec = np.array(np.load(os.path.join(base_path + 'accumulated', str(sess) + '_amp.npy'),
                                          allow_pickle=True))


                reason = np.repeat(amp_sec, fs_new)

                #ax[3].plot(amp_sec, color='magenta', label='iq max amp', linewidth=0.5)
                reason = np.array(np.load(os.path.join(base_path + 'accumulated', str(sess) + '_mamp.npy'),
                                          allow_pickle=True))
                reason = np.repeat(reason, fs_new)
                mean_amp = np.mean(reason)


                #ax[3].plot(reason, color='black', label='iq max amp smooth', linewidth=0.5)

            except:
                print('no amp data for setup', sess)
            try:
                print("motion index")
                reason = np.array(np.load(os.path.join(base_path + 'accumulated', str(sess) + '_diff.npy'),
                                          allow_pickle=True))
                reason = np.repeat(reason, fs_new)

                #ax[3].plot(reason, color='blue', label='motion index', linewidth=0.5)

            except:
                print('no motion index for setup', sess)
        else:


            reason = None
            try:
                print("noise reason")
                reason = np.array(
                    np.load(os.path.join(base_path + 'accumulated', str(sess) + '_save_noise_reason.npy'),
                            allow_pickle=True))
                reason = np.repeat(reason, fs_new)
                for u in np.unique(reason):
                    print(u, 100 * len(reason[reason == u]) / len(reason))
                ax[2].plot(reason, color='magenta', label='noise reason', linewidth=0.5)
                # ax[3].set_yticks(range(-1,5,1))
            except:
                print('no noise reason data for setup', sess)

        comp_stat = None


        computed_hr = None
        all_segments = []
        status_dict = {}
        try:
            comp_ra = np.array(np.load(os.path.join(base_path + 'dynamic', str(sess) + '_ra.npy'),
                                      allow_pickle=True))
            comp_ra = np.repeat(comp_ra, fs_new)
            #ax[3].plot(comp_ra, linewidth=0.5, alpha=1, color='red', label='RA')
        except:
            print("No RA DATA")
        try:
            comp_stat = np.array(np.load(os.path.join(base_path + 'dynamic', str(sess) + '_stat.data'),
                                      allow_pickle=True))

            for uu in np.unique(comp_stat):
                s = pd.DataFrame(comp_stat, columns=['stat'])

                vv = comp_stat == uu
                stat_vec = np.zeros_like(vv)
                stat_vec[vv == True] = 1

                print("Percentage",uu,  100 * (sum(vv) / len(vv)))
                print("comp_stat getting segments for", uu)
                es = getSegments(s, uu)
                #print("OFFLINE", uu,  len(es), es)
                status_dict[uu] = len(es)
                #ax[0].set_title(line'Offline Algo Computed Status')
                for e in es:
                    suu = standardize_key(uu)
                    all_segments.append([to_class[suu], e[0], e[1]])

                    if suu in to_draw.keys() and to_draw[suu]:
                        try:
                            ax[5].axvspan(e[0] * fs_new, e[1] * fs_new, color=cols[suu], alpha=0.3)
                        except:
                            print("plotting", uu, "failed")
                print("done drawing comp_stat")
        except:
            print('no online stat data for setup', sess)





        hr_hq = np.zeros_like(bins)
        try:
            hr_offline_by_sec = np.array(np.load(os.path.join(base_path + 'dynamic', str(sess) + '_hr.npy'),
                                         allow_pickle=True))
            hr_hq = np.array(np.load(os.path.join(base_path + 'dynamic', str(sess) + '_hr_high_quality_indicator.npy'),
                                       allow_pickle=True))


            comp_hr = np.repeat(hr_offline_by_sec, fs_new)
            hr_hq_long = np.repeat(hr_hq, fs_new)
            hr_hq_rate_offline = int(np.round(100 * len(hr_hq_long[hr_hq_long==1])/len(hr_hq_long[hr_hq_long>=0])))

            hr_val_rate_offline = int(np.round(100 * len(comp_hr[comp_hr>0])/len(comp_hr)))



            ax[0].plot(comp_hr, linewidth=1, alpha=0.7, color='red', label='Offline HR: HQ '+str(hr_hq_rate_offline)+'%,  Pos: '+str(hr_val_rate_offline)+'%')


            rr_offline_by_sec = np.array(np.load(os.path.join(base_path + 'dynamic', str(sess) + '_rr.npy'),
                                       allow_pickle=True))
            rr_hq = np.array(np.load(os.path.join(base_path + 'dynamic', str(sess) + '_rr_high_quality_indicator.npy'),
                                     allow_pickle=True))
            comp_rr = np.repeat(rr_offline_by_sec, fs_new)
            rr_hq = np.repeat(rr_hq, fs_new)
            rr_hq_rate_offline = int(np.round(100 * len(rr_hq[rr_hq == 1]) / len(rr_hq[rr_hq >= 0])))
            for s in rr_offline_by_sec:

                if s in range(5,10):
                    offline_10 += 1
                elif s in range(-10, -5):
                    offline_10_neg += 1
                elif s > 32:
                    offline_30 += 1
                elif s < -32:
                    offline_30_neg += 1
            print(offline_10, offline_10_neg, offline_30, offline_30_neg)
            rr_val_rate_offline = int(np.round(100 * len(comp_rr[comp_rr > 0]) / len(comp_rr)))
            ax[0].plot(comp_rr, linewidth=1, alpha=0.7, color='blue', label='Offline RR: HQ '+str(rr_hq_rate_offline)+'%,  Pos: '+str(rr_val_rate_offline)+'%')
            false_lowrr[sess] = {"online_10": online_10, "offline_10": offline_10, "online_30": online_30,
                                 "offline_30": offline_30, "online_10_neg": online_10_neg,
                                 "offline_10_neg": offline_10_neg, "online_30_neg": online_30_neg,
                                 "offline_30_neg": offline_30_neg}

            #ax[0].plot(comp_rr, linewidth=0.65, alpha=0.5, color='green',  label='rr: offline HQ ' + str(rr_hq_rate_offline) + '%,  > -1 ' + str(rr_val_rate_offline) + '%')

        except:
            print('no reject data for setup', sess)

        empty_offline = 0 if 'empty chair' not in status_dict.keys() else status_dict['empty chair']

        # median_hr_offline = np.median(hr_offline_by_sec[hr_offline_by_sec>0])
        # median_hr_online = np.median(hr_online_by_sec[hr_online_by_sec>0])
        #
        # secs_suspected_harmonics_online = len(hr_online_by_sec[hr_online_by_sec > 1.5*median_hr_online])
        # secs_suspected_harmonics_offline = len(hr_offline_by_sec[hr_offline_by_sec > 1.5*median_hr_online])
        # harmony_dict[sess] = [secs_suspected_harmonics_online,secs_suspected_harmonics_offline]
        try:
            ax[0].set_title(str(session) + ", " + str(sess))# + ' '+ db.setup_sn(sess)[0][:2]+'0, SN:' + db.setup_sn(sess)[0][-3:] + ', empty events: ' + str(empty_offline) + ', Offline stat')
            ax[1].set_title(str(session) + ", " + str(sess))# + ' '+ db.setup_sn(sess)[0][:2]+ '0, SN:' + db.setup_sn(sess)[0][-3:] + ', empty events: ' + str(empty_online) + ', Online stat ')
        except:
            continue
            ax[0].set_title(str(session) + ", " + str(sess) + ' ' + db.setup_sn(sess)[0][:2] + '0, SN:' + db.setup_sn(sess)[0][-3:])
            ax[1].set_title(str(session) + ", " + str(sess) + ' ' + db.setup_sn(sess)[0][:2] + '0, SN:' + db.setup_sn(sess)[0][-3:])
        ax[1].legend(loc='upper right', fontsize="7")

        #ax[3].axhline(y=40, color='red', linewidth=0.5)
        #ax[4].axhline(y=0.28, color='red', linewidth=0.5)
        ax[4].axhline(y=0.28, color='red', linewidth=0.5)
        ax[4].axhline(y=0.5, color='red', linewidth=0.5)
        #ax[4].axhline(y=0.4, color='red', linewidth=0.5)
        #ax[2].axhline(y=60, color='magenta', linewidth=0.5)
        #ax[2].axhline(y=20, color='magenta', linewidth=0.5)
        # ax[1].axvline(x=1123650, color='magenta', linewidth=1, alpha=0.5)
        # ax[1].axvline(x=1123900, color='magenta', linewidth=1, alpha=0.5)
        # ax[5].text(1123750, 0, 'suspected harmonics online '+str(secs_suspected_harmonics_online), fontsize=7, alpha=1, color='black')
        # ax[5].text(1123760, 0, 'suspected harmonics 3.6.0 '+str(secs_suspected_harmonics_offline), fontsize=7, alpha=1, color='black')
        #ax[2].set_ylim(-1, 40)
        ax[2].legend(loc='upper right', fontsize="7")
        ax[3].legend(loc='upper right', fontsize="7")
        ax[4].legend(loc='upper right', fontsize="7")
        ax[0].legend(loc='upper right', fontsize="7")
        session = db.session_from_setup(sess)
        outfolder = base_path
        ax[0].axhline(y=10, color='maroon', linewidth=0.25)
        ax[0].axhline(y=30, color='indigo', linewidth=0.25)
        ax[2].axhline(y=10, color='maroon', linewidth=0.25)
        ax[0].axhline(y=0, color='grey', linewidth=0.25)
        ax[2].axhline(y=0, color='grey', linewidth=0.25)
        ax[2].axhline(y=30, color='indigo', linewidth=0.25)
        ax[0].axhline(y=-10, color='maroon', linewidth=0.25)
        ax[0].axhline(y=-30, color='indigo', linewidth=0.25)
        ax[2].axhline(y=-10, color='maroon', linewidth=0.25)
        ax[2].axhline(y=-30, color='indigo', linewidth=0.25)
        slide = ppt.slides.add_slide(ppt.slide_layouts[6])
        image_stream = BytesIO()
        plt.savefig(image_stream)
        slide.shapes.add_picture(image_stream, Cm(0), Cm(0), width=Cm(26), height=Cm(13))


        if shw:
            plt.show()
        fn = '_out.png'
        plt.savefig(outfolder + str(session) + "_" + str(sess) + fn, dpi=300)
        print("saved", outfolder + str(session) + "_" + str(sess) + fn)
        plt.close(fig)

        if db.mysql_db == 'neteera_mirror_db':
            fig4, bx = plt.subplots(2, sharex=False, figsize=(20, 10))
            ml = min(len(amp_sec), len(hr_hq))

            for i in range(ml):
                ind = find_lowest_threshold_index(list(amp_dict.keys()), amp_sec[i])
                amp_dict[list(amp_dict.keys())[ind]].append(hr_hq[i])

            for k,v in amp_dict.items():
                hq_perc = 0 if len(v) == 0 else np.sum(v)/len(v)
                print(k, hq_perc)
                bx[1].scatter(k, hq_perc, s=6)
                bx[1].text(k, hq_perc,len(v), fontsize=5, alpha=0.5)

            bx[1].set_title("HR HQ vs. amplitude")

            ml = min(len(bins), len(hr_hq))
            for i in range(ml):
                bin_dict[bins[i]].append(hr_hq[i])
                if bins[i] == -1 and hr_hq[i] == 1:
                    print("HQ = 1 && bin == -1:",sess, i)
            for k,v in bin_dict.items():
                hq_perc = 0 if len(v) == 0 else np.sum(v) / len(v)
                print(k, hq_perc)
                bx[0].scatter(k, hq_perc, s=6)
                bx[0].text(k, hq_perc,len(v), fontsize=5, alpha=0.5)

            #bx[0].scatter(bins[:ml], hr_hq[:ml], s=4)
            bx[0].set_title("HQ vs. bin #")
            plt.savefig(outfolder + "bins.png", dpi=300)
            print("saved", outfolder + "bins.png")

            print(sess, status_dict)
    for k, v in false_lowrr.items():
        print(k, v)
        sum_10_after += v['offline_10']
        sum_10_before += v['online_10']
        sum_30_after += v['offline_30']
        sum_30_before += v['online_30']

    #print("5-10", sum_10_before, sum_10_after)
    #print(">32", sum_30_before, sum_30_after)
    ppt.save(os.path.join(base_path, 'ppt.pptx'))
    np.save(os.path.join(base_path, 'bin_changes.npy'),binchanges_list, allow_pickle=True)

