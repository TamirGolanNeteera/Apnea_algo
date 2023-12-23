# Copyright (c) 2021 Neteera Technologies Ltd. - Confidential
import os
import sys
sys.path.insert(1, os.getcwd())
# from Tests.Plots.PlotRawDataRadarCPX import*
import argparse
import fnmatch
from os import listdir
from sklearn import preprocessing
import pandas as pd
from Tests.Utils.LoadingAPI import load_reference
from Tests.vsms_db_api import *
from pptx import Presentation
from pptx.util import Inches, Cm
from io import BytesIO
from Tests.NN.create_apnea_count_AHI_data import get_signal_time_index, MB_HQ, delays, getSetupDataLocalDB, getSetupRespirationCloudDBDebug, getSetupRespirationCloudDBDebugWithTimestamp, MB_HQ, count_apneas_in_chunk, getSetupRespirationCloudDBDebugDecimate, getSetupRespirationLocalDBDebug, getSetupRespirationCloudDB, compute_respiration, compute_phase
from Tests.NN.Chest.create_apnea_count_AHI_data_regression_MB_chest import apnea_class
import matplotlib.pyplot as plt
import glob
import scipy.signal as sp
from scipy.signal import spectrogram
from Tests.NN.reiniers_create_data import spectrogram_features
from scipy.signal import butter, filtfilt, lfilter, firwin
import datetime
import matplotlib.pyplot as plt
import matplotlib.dates as mdates

db = DB()
home = '/Neteera/Work/homes/dana.shavit/Research/analysis/'


def find_lowest_threshold_index(A, N):
    # Initialize variables to store the minimum threshold and its index
    min_threshold = float('inf')  # Initialize with positive infinity
    min_index = -1

    # Iterate through the thresholds and find the lowest one greater than N
    for i, threshold in enumerate(A):
        if threshold > N and threshold < min_threshold:
            min_threshold = threshold
            min_index = i

    return min_index

ss_dict = {'N1':1, 'N2':1, 'N3':3, 'R':4, 'W':0, 'invalid':-1}

all_bin_dict = {}
for i in range(-1,15):
    all_bin_dict[i] = []

amp_keys = [250, 500, 750, 1000, 2000, 3000, 4000, 5000, 6000, 7000, 8000, 9000, 10000]
all_amp_dict = {}
for k in amp_keys:
    all_amp_dict[k] = []


def rollavg_convolve_edges(a,n):
    'scipy.convolve, edge handling'
    assert n%2==1
    return sp.convolve(a,np.ones(n,dtype='float'), 'same')/sp.convolve(np.ones(len(a)),np.ones(n), 'same')

base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.6.xx/braintree/'
#base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.6.xx/tamir/'
base_path = '/Neteera/Work/homes/dana.shavit/Research/analysis/3.6.xx/szmc_new_2_selected_sessions/'

setups = ['all']
setups = [111662,111664]

shw=False
if setups == ['all']:
    shw=False
    setups = []
    setup_files = fnmatch.filter(os.listdir(os.path.join(base_path,'dynamic')), '*_hr.npy')

    print("setup_files", setup_files)

    for i, fn in enumerate(setup_files):
        sess = int(fn[0:fn.find('_')])
        setups.append(sess)

print(setups)

def high_pass_bin(signal):

    # Apply high-pass filter with cutoff frequency of 20 Hz
    cutoff_freq = 0.1  # Cutoff frequency (Hz)
    fs = 500
    # Design the high-pass filter
    nyquist_freq = 0.5 * fs
    norm_cutoff_freq = cutoff_freq / nyquist_freq
    b, a = butter(3, norm_cutoff_freq, btype='high', analog=False, output='ba')

    # Apply the filter to the signal
    filtered_signal = filtfilt(b, a, signal)

    return filtered_signal
def get_args() -> argparse.Namespace:
    """ Argument parser

    :return: parsed arguments of the types listed within the function
    :rtype: argparse.Namespace
    """
    parser = argparse.ArgumentParser(description='Process some integers.')
    parser.add_argument('-save_path', metavar='Location', type=str, required=True, help='location of saved data')
    parser.add_argument('--scale', action='store_true', help='Scale data to m=0 s=1')
    parser.add_argument('--show', action='store_true', help='Scale data to m=0 s=1')
    parser.add_argument('--overwrite', action='store_true',  required=False, help='Overwrite existing output')
    parser.add_argument('-chunk', metavar='window', type=int, required=True, help='signal size')
    parser.add_argument('-step', metavar='window', type=int, required=True, help='stride for signal creation')
    return parser.parse_args()

cols = {'empty chair': 'red', 'empty':'red',
                    'low_s': 'yellow',
                    'motio': 'blue',
                    'runni': 'green',
                    'full': 'green',
                    'unkno': 'grey',
                    'warmi': 'magenta',
                    'missi': 'black',
                    'low_r':'black'}

cols_class = {0: 'red', 1: 'yellow', 2: 'blue', 3: 'green', 4: 'green', 5: 'grey', 6: 'olivedrab', 7: 'black', 8:'black'}
to_draw = {'empty':True,
                    'low_s': True,
                    'motio': True,
                    'runni': True,
                    'full': True,
                    'unkno': True,
                    'warmi': True,
                    'missi': False,
                    'low_r':False}
to_class = {'empty':0,
                    'low_s': 1,
                    'motio': 2,
                    'runni': 3,
                    'full': 3,
                    'unkno': 4,
                    'warmi': 5,
                    'missi': 6,
                    'low_r':7}

def standardize_key(curr_key):
    s_key = curr_key.lower()
    s_key = s_key.replace(' ', '_')
    s_key = s_key.replace('-', '_')
    s_key = s_key[:min(5, len(s_key))]
    return s_key


def getSegments(data_frame, seg_val, decimate=False):
    #
    countr = 0
    segments = []
    for k, v in data_frame[data_frame['stat'] == seg_val].groupby((data_frame['stat'] != seg_val).cumsum()):
        if decimate:
            segments.append([int(np.floor(v.index[0])), 1+int(np.floor(v.index[-1]))])
        else:
            segments.append([int(v.index[0]), 1+int(v.index[-1])])
        countr += 1
    return segments


if __name__ == '__main__':
    col = ['gray', 'blue', 'green', 'red', 'yellow', 'magenta', 'cyan']

    print(setups)
    #setups = [109870]
    d = []
    ppt = Presentation()

    title_slide_layout = ppt.slide_layouts[6]
    for i_sess, sess in enumerate(setups):
        db.update_mysql_db(sess)
        radar_sn = db.setup_sn(sess)[0]
        if radar_sn[0] == 1:
            print(sess, radar_sn, "skipping 130 setup")
            continue

        bin_dict = {}
        for i in range(-1, 15):
            bin_dict[i] = []

        amp_keys = [250, 500, 750, 1000, 2000, 3000, 4000, 5000, 6000, 7000, 8000, 9000, 10000]
        amp_dict = {}
        for k in amp_keys:
            amp_dict[k] = []
            
        setup_row = []

        setup_row.append(sess)
        setup_row.append(db.setup_duration(sess))
        #setup_row.append(db.setup_sn(sess)[0])

        ts = None
        use_ts = False
        try:
            if sess < 100000:
                respiration, fs_new, bins, I, Q = getSetupDataLocalDB(sess)

            else:
                respiration, fs_new, bins, I, Q = getSetupRespirationCloudDBDebug(sess)

                print(len(respiration), len(I), len(Q))
                print("LENGTH BINS", len(bins))
                binss = np.repeat(bins, fs_new)
                ax[3].plot(binss, label='bins', linewidth=0.75)
        except:
            print("setup not found, skipping")
        print(fs_new, "HZ")
        ts =  np.arange(0, db.setup_duration(sess), 1/fs_new)
        print("len(ts)", len(ts))
        print("len(I)", len(I))

        session = db.session_from_setup(sess)
        dist = db.setup_distance(sess)
        bin_change = np.diff(bins)

        fig, ax = plt.subplots(6, sharex=True, figsize=(20, 10))

        try:
            ax[5].plot(binss, label='bins', linewidth=0.75)
        except:
            print("no bins")
        bin_change = np.diff(bins)
        ax[5].set_title(str(sess)+' Selected Bins')

        print(len(I))

        ax[0].plot(respiration, label='phase', linewidth=0.25)

        ax[1].plot(I, label='I online', linewidth=0.25)
        ax[1].plot(Q, label='Q online', linewidth=0.25)

        try:
            p = db.setup_dir(sess)
            raw_dir = os.sep.join([p, 'NES_RAW'])
            res_dir = os.sep.join([p, 'NES_RES'])
        except:
            continue


        try:
            onlyfiles = listdir(res_dir)
            print(res_dir)
            vs = [f for f in onlyfiles if 'VS' in f and 'csv' in f]

            if len(onlyfiles) ==1:
                res_dir = os.sep.join([res_dir, onlyfiles[0]])
                onlyfiles = listdir(res_dir)
                vs = [f for f in onlyfiles if 'results' in f and 'csv' in f]

            vs = vs[0]

            df = pd.read_csv(os.sep.join([res_dir, vs]))
            rr = np.repeat(df['rr'].to_numpy(), fs_new)
            rr_hq = df['rr_hq'].to_numpy()
            hr = np.repeat(df['hr'].to_numpy(), fs_new)
            hr_hq = df['hr_hq'].to_numpy()
            print("LENGTH HR HQ", len(hr_hq))

            df['stat'] = df['stat'].fillna('missing')
            #status = df['stat'].to_numpy()
            #print(np.unique(status))
            print(len(df), len(I)/fs_new)
            # if len(df) > int(1.5 * len(I)/fs_new):
            #     print("not ok")
            #     continue

            hr_hq_rate_online = int(np.round(100 * len(hr_hq[hr_hq==1])/len(hr_hq[hr_hq>=0])))
            rr_hq_rate_online = int(np.round(100 * len(rr_hq[rr_hq==1])/len(rr_hq[rr_hq>=0])))

            rr_hq_rate = int(np.round(100 * len(rr_hq[rr_hq==1])/len(rr_hq[rr_hq>=0])))
            hr_val_rate_online = int(np.round(100 * len(hr[hr>0])/len(hr)))
            rr_val_rate_online = int(np.round(100 * len(rr[rr>0])/len(rr)))
            print(hr_hq_rate_online, rr_hq_rate)
            ax[2].set_title('VS')
            ax[2].plot(hr, linewidth=3, alpha=0.3, color='red', label='HR: HQ '+str(hr_hq_rate_online)+'%,  Pos: '+str(hr_val_rate_online)+'%')
            ax[2].plot(rr, linewidth=3, alpha=0.3, color='blue', label='RR: HQ '+str(rr_hq_rate_online)+'%,  Pos: '+str(rr_val_rate_online)+'%')

            print(np.unique(df['stat']))
            empty_online = 0
            for uu in np.unique(df['stat']):
                vv = df['stat'] == uu
                stat_vec = np.zeros_like(vv)
                stat_vec[vv == True] = 1

                print("Percentage", uu, 100 * (sum(vv) / len(vv)))
                print("df['stat'] getting segments for", uu)
                es = getSegments(df, uu)

                #print("ONLINE", uu, es)
                if 'mpty' in uu or 'MPTY' in uu:
                    empty_online = len(es)

                print("success...")
                # ax[1].set_title('Device Online Computed Status')
                for e in es:

                    suu = standardize_key(uu)
                    if suu in to_draw.keys() and to_draw[suu]:
                        try:
                            kkk=0
                            ax[1].axvspan(e[0] * fs_new, e[1] * fs_new, color=cols[suu], alpha=0.3)
                        except:
                            print("plotting", uu, "failed")
                print("done drawing df['stat']")
        except:
            print("no online result data for setup", sess)


        #empty chair reasons
        nt = [] if not db.setup_note(sess) else db.setup_note(sess)
        plot_empty_data = True
        if plot_empty_data:
            reason = None
            try:

                print("loading zero crossing data")
                reason = np.array(np.load(os.path.join(base_path + 'accumulated', str(sess) + '_mzc.npy'),
                                          allow_pickle=True))
                reason = np.repeat(reason, fs_new)

                ax[4].plot(reason, color='black', label='zc %ile smooth', linewidth=0.5)
                mean_noise = np.mean(reason)


            except:
                print('no zc data for setup', sess)
            reason = None
            try:
                print("loading amp data")
                amp_sec = np.array(np.load(os.path.join(base_path + 'accumulated', str(sess) + '_amp.npy'),
                                          allow_pickle=True))


                reason = np.repeat(amp_sec, fs_new)

                ax[3].plot(amp_sec, color='magenta', label='iq max amp', linewidth=0.5)
                reason = np.array(np.load(os.path.join(base_path + 'accumulated', str(sess) + '_mamp.npy'),
                                          allow_pickle=True))
                reason = np.repeat(reason, fs_new)
                mean_amp = np.mean(reason)


                ax[3].plot(reason, color='black', label='iq max amp smooth', linewidth=0.5)

            except:
                print('no amp data for setup', sess)
            try:
                print("motion index")
                reason = np.array(np.load(os.path.join(base_path + 'accumulated', str(sess) + '_diff.npy'),
                                          allow_pickle=True))
                reason = np.repeat(reason, fs_new)

                ax[3].plot(reason, color='blue', label='motion index', linewidth=0.5)

            except:
                print('no motion index for setup', sess)
        else:
            reason = None
            try:
                print("hr == -1 reason")
                reason = np.array(np.load(os.path.join(base_path + 'accumulated', str(sess) + '_quality_fail.npy'),
                                          allow_pickle=True))
                reason = np.repeat(reason, fs_new)
                for u in np.unique(reason):
                    print(u, 100*len(reason[reason == u])/len(reason))
                ax[4].plot(reason, color='magenta', label='LQ reason', linewidth=0.5)
                #ax[3].set_yticks(range(-1, 5, 1))
            except:
                print('no qual fail data for setup', sess)

            reason = None
            try:
                print("noise reason")
                reason = np.array(
                    np.load(os.path.join(base_path + 'accumulated', str(sess) + '_save_noise_reason.npy'),
                            allow_pickle=True))
                reason = np.repeat(reason, fs_new)
                for u in np.unique(reason):
                    print(u, 100 * len(reason[reason == u]) / len(reason))
                ax[2].plot(reason, color='magenta', label='noise reason', linewidth=0.5)
                # ax[3].set_yticks(range(-1,5,1))
            except:
                print('no noise reason data for setup', sess)

        comp_stat = None


        computed_hr = None
        all_segments = []
        status_dict = {}
        try:
            comp_stat = np.array(np.load(os.path.join(base_path + 'dynamic', str(sess) + '_stat.data'),
                                      allow_pickle=True))

            for uu in np.unique(comp_stat):
                s = pd.DataFrame(comp_stat, columns=['stat'])

                vv = comp_stat == uu
                stat_vec = np.zeros_like(vv)
                stat_vec[vv == True] = 1

                print("Percentage",uu,  100 * (sum(vv) / len(vv)))
                print("comp_stat getting segments for", uu)
                es = getSegments(s, uu)
                #print("OFFLINE", uu,  len(es), es)
                status_dict[uu] = len(es)
                #ax[0].set_title('Offline Algo Computed Status')
                for e in es:
                    suu = standardize_key(uu)
                    all_segments.append([to_class[suu], e[0], e[1]])

                    if suu in to_draw.keys() and to_draw[suu]:
                        try:
                            ax[0].axvspan(e[0] * fs_new, e[1] * fs_new, color=cols[suu], alpha=0.3)
                        except:
                            print("plotting", uu, "failed")
                print("done drawing comp_stat")
        except:
            print('no online stat data for setup', sess)

        hr_hq = np.zeros_like(bins)
        try:
            comp_hr = np.array(np.load(os.path.join(base_path + 'dynamic', str(sess) + '_hr.npy'),
                                         allow_pickle=True))
            hr_hq = np.array(np.load(os.path.join(base_path + 'dynamic', str(sess) + '_hr_high_quality_indicator.npy'),
                                       allow_pickle=True))
            comp_hr = np.repeat(comp_hr, fs_new)
            hr_hq_long = np.repeat(hr_hq, fs_new)
            hr_hq_rate_offline = int(np.round(100 * len(hr_hq_long[hr_hq_long==1])/len(hr_hq_long[hr_hq_long>=0])))

            hr_val_rate_offline = int(np.round(100 * len(comp_hr[comp_hr>0])/len(comp_hr)))


            setup_row.append(hr_hq_rate_online)
            setup_row.append(hr_hq_rate_offline)

            setup_row.append(hr_val_rate_online)
            setup_row.append(hr_val_rate_offline)


            ax[2].plot(comp_hr, linewidth=1, alpha=1, color='red', label='Offline HR: HQ '+str(hr_hq_rate_offline)+'%,  Pos: '+str(hr_val_rate_offline)+'%')


            comp_rr = np.array(np.load(os.path.join(base_path + 'dynamic', str(sess) + '_rr.npy'),
                                       allow_pickle=True))
            rr_hq = np.array(np.load(os.path.join(base_path + 'dynamic', str(sess) + '_rr_high_quality_indicator.npy'),
                                     allow_pickle=True))
            comp_rr = np.repeat(comp_rr, fs_new)
            rr_hq = np.repeat(rr_hq, fs_new)
            rr_hq_rate_offline = int(np.round(100 * len(rr_hq[rr_hq == 1]) / len(rr_hq[rr_hq >= 0])))

            rr_val_rate_offline = int(np.round(100 * len(comp_rr[comp_rr > 0]) / len(comp_rr)))
            ax[2].plot(comp_rr, linewidth=1, alpha=1, color='blue', label='Offline RR: HQ '+str(rr_hq_rate_offline)+'%,  Pos: '+str(rr_val_rate_offline)+'%')

            # ax[2].plot(comp_rr, linewidth=0.65, alpha=0.5, color='green',
            #            label='rr: offline HQ ' + str(rr_hq_rate_offline) + '%,  > -1 ' + str(rr_val_rate_offline) + '%')
        except:
            print('no reject data for setup', sess)

        empty_offline = 0 if 'empty chair' not in status_dict.keys() else status_dict['empty chair']

        try:
            ax[0].set_title(str(session) + ", " + str(sess) + ' '+ db.setup_sn(sess)[0][:2]+'0, SN:' + db.setup_sn(sess)[0][-3:] + ', empty events: ' + str(empty_offline) + ', Offline stat')
            ax[1].set_title(str(session) + ", " + str(sess) + ' '+ db.setup_sn(sess)[0][:2]+ '0, SN:' + db.setup_sn(sess)[0][-3:] + ', empty events: ' + str(empty_online) + ', Online stat ')
        except:
            ax[0].set_title(str(session) + ", " + str(sess) + ' ' + db.setup_sn(sess)[0][:2] + '0, SN:' + db.setup_sn(sess)[0][-3:])
            ax[1].set_title(str(session) + ", " + str(sess) + ' ' + db.setup_sn(sess)[0][:2] + '0, SN:' + db.setup_sn(sess)[0][-3:])
        ax[1].legend(loc='upper right', fontsize="7")
        ax[3].axhline(y=250, color='red', linewidth=0.5)
        ax[3].axhline(y=1000, color='red', linewidth=0.5)
        ax[3].axhline(y=2500, color='red', linewidth=0.5)
        #ax[3].axhline(y=5000, color='red', linewidth=0.5)
        ax[3].axhline(y=5000, color='red', linewidth=0.5)
        #ax[3].axhline(y=20, color='red', linewidth=0.5)
        #ax[3].axhline(y=40, color='red', linewidth=0.5)
        #ax[4].axhline(y=0.28, color='red', linewidth=0.5)
        ax[4].axhline(y=0.75, color='red', linewidth=0.5)
        ax[4].axhline(y=0.5, color='red', linewidth=0.5)
        #ax[4].axhline(y=0.4, color='red', linewidth=0.5)
        ax[2].axhline(y=60, color='magenta', linewidth=0.5)
        ax[2].axhline(y=20, color='magenta', linewidth=0.5)
        #ax[2].set_ylim(-1, 40)
        ax[2].legend(loc='upper right', fontsize="7")
        ax[3].legend(loc='upper right', fontsize="7")
        ax[4].legend(loc='upper right', fontsize="7")
        ax[0].legend(loc='upper right', fontsize="7")
        session = db.session_from_setup(sess)
        outfolder = base_path

        slide = ppt.slides.add_slide(ppt.slide_layouts[6])
        image_stream = BytesIO()
        plt.savefig(image_stream)
        slide.shapes.add_picture(image_stream, Cm(0), Cm(0), width=Cm(26), height=Cm(13))
        if shw:
            plt.show()
        fn = '_out.png'
        plt.savefig(outfolder + str(session) + "_" + str(sess) + fn, dpi=300)
        print("saved", outfolder + str(session) + "_" + str(sess) + fn)
        plt.close(fig)
        fig4, bx = plt.subplots(2, sharex=False, figsize=(20, 10))
        ml = min(len(amp_sec), len(hr_hq))

        for i in range(ml):
            ind = find_lowest_threshold_index(list(amp_dict.keys()), amp_sec[i])
            amp_dict[list(amp_dict.keys())[ind]].append(hr_hq[i])
            all_amp_dict[list(all_amp_dict.keys())[ind]].append(hr_hq[i])

        for k,v in amp_dict.items():
            hq_perc = 0 if len(v) == 0 else np.sum(v)/len(v)
            print(k, hq_perc)
            bx[1].scatter(k, hq_perc, s=6)
            bx[1].text(k, hq_perc,len(v), fontsize=5, alpha=0.5)

        bx[1].set_title(str(sess)+" HR HQ vs. amplitude " + str(hr_hq_rate_online)+" %HR HQ")

        ml = min(len(bins), len(hr_hq))
        for i in range(ml):
            bin_dict[bins[i]].append(hr_hq[i])
            all_bin_dict[bins[i]].append(hr_hq[i])
            if bins[i] == -1 and hr_hq[i] == 1:
                print("HQ = 1 && bin == -1:",sess, i)
        for k,v in bin_dict.items():
            hq_perc = 0 if len(v) == 0 else np.sum(v) / len(v)
            print(k, hq_perc)
            bx[0].scatter(k, hq_perc, s=6)
            bx[0].text(k, hq_perc,len(v), fontsize=5, alpha=0.5)

        #bx[0].scatter(bins[:ml], hr_hq[:ml], s=4)
        bx[0].set_title(str(sess)+" HQ vs. bin # " + str(hr_hq_rate_online)+" %HR HQ")
        plt.savefig(outfolder + str(sess)+"_bins.png", dpi=300)
        if shw:
            plt.show()
        plt.close()

        print(sess, status_dict)
    ppt.save(os.path.join(base_path, 'ppt.pptx'))
    plt.close()
    fig3, cx = plt.subplots(2, sharex=False, figsize=(20, 10))
    for k,v in all_amp_dict.items():
        hq_perc = 0 if len(v) == 0 else np.sum(v)/len(v)
        print(k, hq_perc)
        cx[1].scatter(k, hq_perc, s=6)
        cx[1].text(k, hq_perc,len(v), fontsize=5, alpha=0.5)

        cx[1].set_title("ALL HR HQ vs. amplitude")

    for k,v in all_bin_dict.items():
        hq_perc = 0 if len(v) == 0 else np.sum(v) / len(v)
        print(k, hq_perc)
        cx[0].scatter(k, hq_perc, s=6)
        cx[0].text(k, hq_perc,len(v), fontsize=5, alpha=0.5)

        #cx[0].scatter(bins[:ml], hr_hq[:ml], s=4)
        cx[0].set_title("ALL HQ vs. bin no.")
    plt.savefig(outfolder + "all_bins.png", dpi=300)
    plt.close()
    print("saved", outfolder + "all_bins.png")